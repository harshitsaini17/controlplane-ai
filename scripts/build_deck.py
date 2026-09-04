"""Build b24bb1029_ControlPlane_Business_Proposal.pptx from the brand template.

Every figure on every slide comes from b24bb1029_ControlPlane_Business_Proposal_FINAL.md
and carries its [B-nn] / [M-nn] provenance tag. Diagrams are native PowerPoint shapes
(see deck_diagrams.py) -- no rasterised pictures.

Run:  .venv/bin/python scripts/build_deck.py
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import deck_diagrams as D
from deck_lib import (ACCENT, BLACK, BODY_TOP, CW, DEEP, GREY_BG, MUTED, ML, RULE, SRC_Y,
                      TINT1, TINT2, WHITE, Deck, bold_italic_from_facename, bullets,
                      delete_shape, divider, find_shape, footer, force_arial, hline,
                      kpi_tile, lead, para, rect, reorder_and_prune, set_text,
                      slide_title, source_line, style_run, table, textbox)

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "assets" / "AIC_Talent-Brand_PPT-Template.pptx"
OUT = ROOT / "submission" / "b24bb1029_ControlPlane_Business_Proposal.pptx"

REPRO_EVAL = "Source: reports/eval_report.md · reproduce: python -m eval.run_all"
REPRO_LAT = "Source: reports/latency_report.md · reproduce: python -m eval.bench_latency --check"


# =============================================================== template-slide edits
def fix_title_slide(slide):
    """Template cover: keep the brand artwork, set our product line."""
    tb = find_shape(slide, "TextBox 15")
    set_text(tb, ["ControlPlane"], sizes=[Pt(72)], colors=[WHITE], bolds=[True],
             aligns=[PP_ALIGN.CENTER])
    tb.left, tb.top, tb.width = Inches(2.50), Inches(6.22), Inches(15.0)
    tb.height = Inches(1.32)

    sub = textbox(slide, 2.50, 7.60, 15.0, 1.52)
    para(sub.text_frame,
         "The real-time oversight gateway for enterprise AI", 22, color=WHITE,
         align=PP_ALIGN.CENTER, first=True)
    para(sub.text_frame,
         "One URL change. Every response checked across performance, cost and "
         "responsibility, and one accountable verdict before delivery.",
         13, color=TINT1, align=PP_ALIGN.CENTER, space_before=8, line=1.10)

    team = textbox(slide, 2.50, 9.24, 15.0, 0.42)
    para(team.text_frame,
         "Team b24bb1029 · Priyanshu Pandey · Harshit Saini · Jayant Soni · "
         "Indian Institute of Technology Jodhpur", 12, color=WHITE,
         align=PP_ALIGN.CENTER, first=True)

    # the template ships a 2025 copyright on the cover; the deck standard is 2026
    for sh in slide.shapes:
        if sh.shape_type == 6:
            for sub_sh in sh.shapes:
                if sub_sh.has_text_frame and "2025 Accenture" in sub_sh.text_frame.text:
                    set_text(sub_sh, ["Copyright © 2026 Accenture. All rights reserved."],
                             sizes=[Pt(12)], colors=[MUTED])
                if sub_sh.has_text_frame and sub_sh.text_frame.text.strip() == "‹#›":
                    set_text(sub_sh, ["1"], sizes=[Pt(12)], colors=[MUTED],
                             aligns=[PP_ALIGN.RIGHT])


MEMBERS = [
    # (name box, detail box, name, roll, stream, leader?)
    ("TextBox 12", "TextBox 30", "Priyanshu Pandey", "B24BB1029", "Bioengineering", True),
    ("TextBox 21", "TextBox 31", "Harshit Saini", "B24CS1031", "Computer Science", False),
    ("TextBox 27", "TextBox 32", "Jayant Soni", "B24CM1033",
     "Artificial Intelligence & Data Science", False),
]


def fill_team_slide(slide):
    """Populate the template-mandated team slide. Layout and photo frames untouched."""
    tbl = [sh for sh in slide.shapes if sh.has_table][0].table
    for ci, txt in ((0, "TEAM NAME:  b24bb1029"),
                    (1, "Indian Institute of Technology Jodhpur")):
        cell = tbl.cell(0, ci)
        cell.text_frame.word_wrap = True
        for p in list(cell.text_frame.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        p0 = cell.text_frame.paragraphs[0]
        for r in list(p0.runs):
            r._r.getparent().remove(r._r)
        style_run(p0.add_run(), Pt(20), bold=True, color=BLACK).text = txt

    def named(name):
        for sh in slide.shapes:
            if sh.shape_type == 6:
                for s2 in sh.shapes:
                    if s2.name == name:
                        return s2
            if sh.name == name:
                return sh
        return None

    for name_box, detail_box, person, roll, stream, leader in MEMBERS:
        # The roll number used to be a second line inside the name plate. That plate is
        # picture-filled with the template's hill graphic, so grey-on-olive rendered the
        # roll effectively unreadable (and purple on green is the same contrast failure).
        # It moves to the white detail block below, where it measures 21:1 / 5.3:1.
        # The plate is BOTTOM-anchored, so the line cannot simply be dropped: without a
        # spacer the name slides down off the light sky band and into the green.
        nb = named(name_box)
        set_text(nb, [person, " "], sizes=[Pt(26), Pt(14)],
                 colors=[BLACK, BLACK], bolds=[True, False])
        db = named(detail_box)
        set_text(db,
                 [f"{roll}{'  ·  TEAM LEADER' if leader else ''}",
                  "College: Indian Institute of Technology Jodhpur",
                  f"Stream: {stream}", "Year of graduation: 2028"],
                 sizes=[Pt(14), Pt(15), Pt(15), Pt(15)],
                 colors=[ACCENT if leader else BLACK, BLACK, BLACK, BLACK],
                 bolds=[True, False, False, False])
        db.width = Inches(4.60)
        db.height = Inches(1.00)   # 4 lines need 0.98in; bottom 6.18 clears the 6.20 rule

    # The template labels only ONE of its three photo frames "Photo"; on a judge-facing
    # mandatory slide that asymmetry reads as unfinished. Label the other two to match
    # (same 27pt white Arial, centred). Appended last, so z-order cannot hide them.
    for px, py in ((1.40, 4.11), (12.12, 4.11)):
        pb = textbox(slide, px, py, 2.37, 2.41, anchor=MSO_ANCHOR.MIDDLE)
        para(pb.text_frame, "Photo", Pt(27), color=WHITE, align=PP_ALIGN.CENTER,
             first=True)

    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "3":
            set_text(sh, ["2"], sizes=[Pt(12)], colors=[MUTED], aligns=[PP_ALIGN.RIGHT])
    # the template team slide carries a page number but no copyright line
    cb = textbox(slide, 14.82, 10.55, 3.93, 0.24)
    para(cb.text_frame, "Copyright © 2026 Accenture. All rights reserved.", Pt(12),
         color=MUTED, first=True)


def fix_thankyou_slide(slide, number):
    for sh in slide.shapes:
        if sh.shape_type == 6:
            for s2 in sh.shapes:
                if s2.has_text_frame and s2.text_frame.text.strip() == "‹#›":
                    set_text(s2, [str(number)], sizes=[Pt(12)], colors=[MUTED],
                             aligns=[PP_ALIGN.RIGHT])
    tb = textbox(slide, 4.55, 7.55, 11.0, 1.00)
    para(tb.text_frame, "ControlPlane: oversight for AI, built under oversight.", 17,
         bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
    para(tb.text_frame,
         "Team b24bb1029 · Indian Institute of Technology Jodhpur · 30 August 2026",
         12, color=TINT1, align=PP_ALIGN.CENTER, space_before=8)


# ================================================================== content slides
def s_exec_summary(s):
    slide_title(s, "Enterprises adopted AI faster than they can supervise it",
                kicker="Executive summary")
    lead(s, "ControlPlane is a reverse-proxy oversight gateway. Applications change one "
            "base URL; every response is then scored across three planes and converged "
            "into exactly one verdict (Pass, Edit, Block or Escalate) before delivery.")
    bullets(s, ML, 3.20, 11.30, [
        ("72% of organisations ", "use generative AI; 80%+ of pilots never reach production [M-03]"),
        ("EU AI Act ", "exposure reaches €35M or 7% of global turnover [M-01]"),
        ("DPDP Act ", "exposure reaches ₹250 crore per instance [M-02]"),
        ("The governing policy is a YAML file ", "per business pipeline"),
        ("Three pipelines, three verdicts ", "on identical content, zero code difference [B-01]"),
        ("No LLM in the decision path: ", "the verdict layer is deterministic"),
    ], pt=12.5, gap=0.52)

    box = rect(s, 12.55, 3.14, 6.69, 3.00, GREY_BG, stroke=RULE)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.20)
    para(tf, "THE ASK", 10.5, bold=True, color=ACCENT, first=True)
    para(tf, "INR 1.6 crore (~US$190,000)", 21, bold=True, color=DEEP, space_before=6)
    para(tf, "Indicative pre-seed for 12 months, to convert this prototype into three "
             "design-partner pilots within two quarters.", 12, color=BLACK,
         space_before=6, line=1.06)
    para(tf, "[model] · §10.4 of the proposal", 10, color=MUTED, space_before=8)

    tiles = [("840/840", "policy-engine verdicts correct across the frozen corpus × 3 policies", "[B-01]"),
             ("1.000", "tier-1 PII precision, zero false alarms", "[B-02]"),
             ("0.133 ms", "tier-1 PII detector P99 against a 2 ms budget", "[B-05]"),
             ("39/39", "fault-injection assertions across 5 quiet-host runs", "[B-07]"),
             ("1,200+", "tests green on two machines and CI", "[B-09]")]
    w, gap = 3.45, 0.32
    x = ML
    for big, cap, tag in tiles:
        kpi_tile(s, x, 6.70, w, 2.30, big, cap, tag=tag, fill=TINT2, big_pt=25)
        x += w + gap
    source_line(s, REPRO_EVAL + "  ·  " + REPRO_LAT)


def s_problem(s):
    slide_title(s, "Three different failures. One identical root cause.",
                kicker="The problem")
    lead(s, "A bank runs three AI assistants. Three teams discover three incidents, "
            "weeks later, from three different systems.")
    cards = [
        ("WRONG", "The customer chatbot confidently states a loan approval that never "
                  "happened.", "Performance plane"),
        ("LEAKY", "The internal HR assistant repeats an employee's phone number into a "
                  "chat transcript.", "Responsibility plane"),
        ("EXPENSIVE", "The advisory tool loops on a malformed integration and burns "
                      "inference budget for a week.", "Cost plane"),
    ]
    w, gap = 5.92, 0.36
    x = ML
    for title, body, plane in cards:
        card = rect(s, x, 3.30, w, 2.70, WHITE, stroke=RULE)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.26)
        tf.margin_top = Inches(0.22)
        para(tf, title, 26, bold=True, color=ACCENT, first=True)
        para(tf, body, 13, color=BLACK, space_before=8, line=1.08)
        para(tf, plane.upper(), 10.5, bold=True, color=MUTED, space_before=10)
        x += w + gap

    band = rect(s, ML, 6.68, CW, 1.62, DEEP)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.36)
    para(tf, "Nothing is supervising the AI at the moment it speaks.", 24, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, first=True)
    tb = textbox(s, ML, 8.62, CW, 0.80)
    para(tb.text_frame,
         "Observability platforms report after delivery. Guardrail engines filter one "
         "category of harm. Gateways route traffic. None makes one accountable, "
         "policy-governed decision per response across all three risk planes.", 12.5,
         color=BLACK, first=True, line=1.10)
    source_line(s, "Source: proposal §3.1 · McKinsey State of AI 2025 and Gartner "
                   "survey series [M-03]")


def s_problem_evidence(s):
    slide_title(s, "The market is a funnel that jams at production",
                kicker="Problem evidence")
    rows = [
        ["Evidence", "Figure", "Source"],
        ["Enterprises using genAI in at least one function", "72%, up from 33% the prior year", "[M-03]"],
        ["Enterprise genAI pilots that never reach production", "over 80%", "[M-03]"],
        ["Enterprises reporting >5% EBIT impact from genAI", "6%", "[M-03]"],
        ["Average production genAI use cases per large enterprise", "5.0, up 101% year over year", "[M-03]"],
        ["Documented, largely unvalidated use cases per enterprise", "~211", "[M-03]"],
        ["Advanced AI projects cancelled for lack of governance", "40%", "[M-12]"],
        ["AI systems classified high-risk / unclassified grey area", "18% / 40%", "[M-12]"],
        ["Estimated recurring EU-wide AI compliance cost", "€3.3 billion annually", "[M-12]"],
        ["Unshielded LLM endpoints vulnerable to injection", "up to 80%+", "[M-05]"],
    ]
    table(s, ML, 2.62, 12.05, rows, [6.2, 4.0, 1.1], pt=11.5, header_pt=11.5,
          row_h=0.60, header_h=0.52)

    box = rect(s, 13.00, 2.62, 6.24, 3.28, GREY_BG, stroke=RULE)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = Inches(0.22)
    para(tf, "READ THE FUNNEL", 10.5, bold=True, color=ACCENT, first=True)
    para(tf, "~211 ideas", 20, bold=True, color=DEEP, space_before=8)
    para(tf, "5.0 live use cases", 20, bold=True, color=DEEP, space_before=4)
    para(tf, "80%+ of pilots dead", 20, bold=True, color=ACCENT, space_before=4)
    para(tf, "Governance is not the only reason pilots stall, but it is a named one, "
             "and the reason a 40% cancellation statistic attaches to directly.", 11.5,
         color=BLACK, space_before=8, line=1.06)

    box2 = rect(s, 13.00, 5.96, 6.24, 2.36, DEEP)
    tf2 = box2.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    tf2.margin_left = tf2.margin_right = Inches(0.26)
    tf2.margin_top = Inches(0.22)
    para(tf2, "THE SECURITY DIMENSION", 10.5, bold=True, color=TINT1, first=True)
    para(tf2, "OWASP ranks prompt injection #1 (LLM01) and sensitive-information "
              "disclosure #6 (LLM06).", 12.5, color=WHITE, space_before=8, line=1.08)
    para(tf2, "Both are response-path problems. Neither is solvable by an offline "
              "evaluation report: by the time it runs, the response has shipped.",
         11.5, color=TINT1, space_before=6, line=1.06)
    source_line(s, "Sources: McKinsey State of AI 2025 / Gartner [M-03] · appliedAI and "
                   "DIGITALEUROPE compliance study, Jul 2026 [M-12] · OWASP GenAI LLM "
                   "Top 10, Aug 2026 [M-05]")


def s_regulatory(s):
    slide_title(s, "The regulatory clock has already started",
                kicker="Regulatory forcing function")
    rows = [
        ["Regime", "Exposure", "Status", "Src"],
        ["EU AI Act\n(Reg. (EU) 2024/1689)",
         "€35M / 7% turnover: prohibited practices\n€15M / 3%: high-risk, GPAI, transparency\n€7.5M / 1%: misrepresentation",
         "In force. Prohibited-practice ban since Feb 2025; GPAI obligations and "
         "administrative fines active since Aug 2025; Annex III high-risk deferred to "
         "Dec 2027 by Reg. (EU) 2026/1744", "[M-01]"],
        ["India DPDP Act 2023",
         "Up to ₹250 crore per instance: failure of reasonable security safeguards\n"
         "₹200 crore: breach-notification failure\n₹150 crore: significant-data-fiduciary duties",
         "Rules notified; Data Protection Board adjudication framework operational",
         "[M-02]"],
        ["Sectoral regimes\n(RBI/SEBI, HIPAA, model risk)",
         "Licence, audit and reputational exposure", "Ongoing", "[M-12]"],
    ]
    table(s, ML, 2.55, CW, rows, [3.1, 6.2, 7.0, 0.9], pt=11.5, header_pt=11.5,
          row_h=1.85, header_h=0.50)

    y = 8.03
    for title, body in [
        ("Obligations vary by use case and jurisdiction.",
         "The same model serving a chatbot and a credit-decision tool sits in two risk "
         "classes. Any system that hard-codes one global policy is already wrong."),
        ("Compliance is becoming a procurement requirement.",
         "Buyers now demand automated logging, dataset lineage and risk-tier evidence "
         "in vendor RFPs."),
    ]:
        card = rect(s, ML, y, CW, 0.78, GREY_BG, stroke=RULE)
        tf = card.text_frame
        tf.margin_left = Inches(0.26)
        p = para(tf, "", 12.5, first=True)
        p.runs[0].text = ""
        style_run(p.add_run(), 12.5, bold=True, color=DEEP).text = title + "  "
        style_run(p.add_run(), 12.5, color=BLACK).text = body
        y += 0.86
    source_line(s, "Sources: eur-lex.europa.eu [M-01] · MeitY / PIB DPDP Rules 2025 "
                   "notification [M-02] · appliedAI / DIGITALEUROPE [M-12]",
                y=SRC_Y - 0.06)


def s_solution_overview(s):
    slide_title(s, "One URL change puts every response under supervision",
                kicker="Solution overview")
    lead(s, "ControlPlane is a reverse proxy between an application and its LLM "
            "provider. The application's code, prompts and provider relationship are "
            "untouched.")
    steps = [
        ("1 · Ingress", "Resolves which business pipeline this request belongs to, and "
                        "loads that pipeline's policy version."),
        ("2 · Input lane", "Pattern detectors and an injection classifier run before the "
                           "provider is called. The provider never receives the raw PII value."),
        ("3 · Dispatch", "Any OpenAI-compatible provider, cloud or fully local, with "
                         "tiered routing and a boot-time accounting canary."),
        ("4 · Output lane", "Streamed responses are buffered and checked sentence by "
                            "sentence. High-stakes pipelines buffer fully."),
        ("5 · Policy engine", "All signals converge into exactly one verdict, by "
                              "documented severity order."),
        ("6 · Audit and review", "Append-only record: categories, never raw values. "
                                 "Escalations enter a human review queue with full lineage."),
    ]
    y = 3.24
    for i, (head, body) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = ML + col * (CW / 2 + 0.20)
        yy = y + row * 1.52
        card = rect(s, x, yy, CW / 2 - 0.20, 1.32, WHITE, stroke=RULE)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.24)
        tf.margin_top = Inches(0.16)
        para(tf, head, 13.5, bold=True, color=DEEP, first=True)
        para(tf, body, 11.5, color=BLACK, space_before=4, line=1.06)

    band = rect(s, ML, 8.30, CW, 1.06, ACCENT)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.34)
    para(tf, "No LLM sits in the decision path. Same signals plus same policy version "
             "always produce the same verdict, which is what makes it usable as "
             "compliance evidence.", 14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, first=True, line=1.06)
    source_line(s, "Source: proposal §4.1 and §5.1 · architecture in docs/02-architecture.md")


def s_architecture(s):
    slide_title(s, "The request lifecycle", kicker="Architecture")
    lead(s, "Six stages, one decision point. The policy engine is the only place a "
            "verdict is issued.")
    D.lifecycle(s)
    source_line(s, "Source: docs/02-architecture.md and docs/04-policy-and-detection-spec.md "
                   "· verdict conformance 840/840 · reproduce: python -m eval.run_all [B-01]")


def s_one_signal(s):
    slide_title(s, "One signal, three business consequences",
                kicker="Why convergence matters")
    D.one_signal_three(s)
    source_line(s, "Source: proposal §5.2 · multi-label signal contract in "
                   "docs/04-policy-and-detection-spec.md")


def s_four_verdicts(s):
    slide_title(s, "The four verdicts", kicker="Decision layer")
    lead(s, "Every response leaves the gateway with exactly one verdict, resolved by "
            "documented severity order, never by filter order.")
    D.four_verdicts(s)
    source_line(s, "Source: verdict resolution 840/840 across the frozen corpus × 3 "
                   "policies · reproduce: python -m eval.run_all [B-01]")


def s_signature(s):
    slide_title(s, "Same content. Three policies. Three verdicts.",
                kicker="The signature behaviour")
    lead(s, "This is the thesis of the product: business context, not model "
            "behaviour, decides what happens to a response.", w=5.62)
    D.three_verdict_fork(s)
    source_line(s, "Source: reports/eval_report.md · 840/840 verdicts across 280 frozen "
                   "cases × 3 policies, enforced by a test that fails if any use-case "
                   "name reaches executable code · reproduce: python -m eval.run_all [B-01]")


def s_escalation(s):
    slide_title(s, "How AI enables it: a tiered escalation ladder",
                kicker="Methodology")
    lead(s, "Latency is a hard constraint: published guidance puts acceptable inline "
            "guardrail overhead at under 20–50 ms [M-11]. So the watching side is "
            "deliberately tiered.")
    D.escalation_ladder(s)
    source_line(s, "Rungs 0–2 measured (reports/latency_report.md); Rungs 3–4 are "
                   "policy-governed escalation. Cascade routing fraction unmeasured: "
                   "SL-10. Guidance band: Lakera / AWS Bedrock [M-11]")


def s_ships(s):
    slide_title(s, "What ships today, and what does not", kicker="Scope of work")
    lb = textbox(s, ML, 2.42, 9.10, 0.42)
    para(lb.text_frame, "SHIPPED IN THE PROTOTYPE (P0, COMPLETE)", 11.5, bold=True,
         color=DEEP, first=True)
    shipped = [
        "Reverse-proxy gateway, OpenAI-compatible, streaming and non-streaming",
        "Three production-shaped pipeline policies (support / HR / regulated advisory)",
        "Input-lane pre-dispatch PII redaction and injection blocking",
        "Sentence-buffered output interception",
        "Six live detectors across the performance and responsibility planes",
        "Deterministic policy engine, four verdicts, multi-label convergence",
        "Per-pipeline fail-open / fail-closed failure posture",
        "Append-only audit store; human review queue with decision lineage",
        "Governance console, landing page, live test-chat interface",
        "Frozen 280-case evaluation corpus and full measurement harness",
        "Scripted eight-beat demonstration with offline replay",
        "CI on two Python versions; two-machine runtime certification",
    ]
    bullets(s, ML, 2.98, 9.10, shipped, pt=11.5, gap=0.455, marker=ACCENT)

    rb = textbox(s, 10.30, 2.42, 8.94, 0.42)
    para(rb.text_frame, "NOT BUILT: PUBLISHED AS STANDING LIMITATIONS", 11.5, bold=True,
         color=ACCENT, first=True)
    rows = [
        ["Gap", "Status", "Plan"],
        ["Cost-plane enforcement detectors",
         "Policy schema and tiered routing ship; enforcement detectors are stubs", "P1"],
        ["Override-driven threshold automation",
         "Review decisions captured with full lineage; the proposal script is a stub", "P1"],
        ["Self-consistency scorer",
         "Specified, unimplemented; the regulated pipeline's quality plane uses grounding", "P1"],
        ["Calibrated thresholds",
         "Calibration ran and inverted on our corpus; seeded thresholds ship, labelled", "P1"],
        ["Conversation-level tracking, deep-audit lane", "Specified, unbuilt", "P2"],
    ]
    table(s, 10.30, 2.98, 8.94, rows, [3.3, 4.9, 0.75], pt=11, header_pt=11,
          row_h=1.02, header_h=0.48)
    nb = textbox(s, 10.30, 8.72, 8.94, 0.70)
    para(nb.text_frame,
         "Nine standing limitations are published in the repository rather than omitted "
         "here. The roadmap in §11 is that register, in priority order.", 11.5,
         italic=True, color=MUTED, first=True, line=1.06)
    source_line(s, "Source: proposal §6.1–6.2 · docs/08-open-questions.md · "
                   "Appendix B (standing limitations SL-1 … SL-9)")


def s_market(s):
    slide_title(s, "Two analyst estimates, and why they differ", kicker="Market size")
    D.market_funnel(s)
    source_line(s, "Sources: MarketsandMarkets, Jan 2025 and Grand View Research, "
                   "Sep 2024 [M-04] · serviceable segment: NASSCOM / MeitY GCC "
                   "reporting, Nov 2025 [M-09] · SOM is [model]")


def s_competitive(s):
    slide_title(s, "Nobody occupies the converged, pre-delivery quadrant",
                kicker="Competitive landscape")
    D.positioning_map(s)
    source_line(s, "Sources: Tracxn, Crunchbase, and AWS / Microsoft / NVIDIA official "
                   "documentation, Feb 2026 [M-07] · Presidio baseline arXiv:2608.02616v1 "
                   "[M-08]")


def s_retraction(s):
    slide_title(s, "The claim we withdrew when the evidence contradicted it",
                kicker="Positioning integrity")
    band = rect(s, ML, 2.42, CW, 2.06, DEEP)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.34)
    tf.margin_top = Inches(0.20)
    para(tf, "RETRACTED", 11, bold=True, color=ACCENT, first=True)
    para(tf, "An earlier draft claimed that per-use-case policy configuration and "
             "configurable fail-open/fail-closed posture were unique to ControlPlane.",
         14, bold=True, color=WHITE, space_before=6, line=1.06)
    para(tf, "Primary-source research showed that claim to be false: AWS Bedrock "
             "Guardrails, Azure AI Content Safety and NVIDIA NeMo Guardrails all "
             "document both capabilities [M-07]. We removed the claim, three weeks "
             "before anyone would have checked.", 12, color=TINT1, space_before=6,
         line=1.06)

    hb = textbox(s, ML, 4.72, CW, 0.40)
    para(hb.text_frame, "THE FIVE DIFFERENTIATORS THAT SURVIVE THE EVIDENCE", 11.5,
         bold=True, color=DEEP, first=True)
    items = [
        ("Convergence across three planes into one verdict.",
         "The hyperscaler engines are content-safety systems. None converges quality, "
         "cost and responsibility into a single accountable decision per response."),
        ("Multi-label overlapping risk.",
         "One signal carrying both a hallucination and a privacy label, resolved by the "
         "engine rather than by filter order."),
        ("Provider and deployment neutrality.",
         "Cross-provider by construction, including fully local and air-gapped operation "
         "which is structurally unavailable from a cloud-native guardrail service."),
        ("Evidence as a product surface.",
         "Frozen corpus, blind-first measurement, published misses, claim-to-command map. "
         "No competitor in our research publishes an auditable accuracy claim about its "
         "own detectors."),
        ("A process record that cannot be reconstructed retroactively.",
         "The adjudication ledger accrues in the order it was created."),
    ]
    y = 5.18
    for i, (head, body) in enumerate(items, 1):
        card = rect(s, ML, y, CW, 0.82, WHITE if i % 2 else GREY_BG, stroke=RULE)
        tf = card.text_frame
        tf.margin_left = tf.margin_right = Inches(0.24)
        p = para(tf, "", 12, first=True)
        p.runs[0].text = ""
        style_run(p.add_run(), 12, bold=True, color=ACCENT).text = f"{i}.  "
        style_run(p.add_run(), 12, bold=True, color=DEEP).text = head + "  "
        style_run(p.add_run(), 12, color=BLACK).text = body
        y += 0.88

    tb = textbox(s, ML, 9.62, CW, 0.40)
    para(tb.text_frame, "We would rather present five defensible differentiators than "
                        "six, one of which fails on inspection.", 13, bold=True,
         color=DEEP, first=True)
    source_line(s, "Source: proposal §6.4 · competitor capability evidence [M-07]",
                y=SRC_Y + 0.14)


def s_users(s):
    slide_title(s, "Who feels the pain, and who signs", kicker="Target users and buyers")
    rows = [
        ["Persona", "Pain today", "What ControlPlane provides", "Evidence of pain"],
        ["Head of Platform /\nAI Engineering",
         "Every new use case re-solves oversight; N apps × M providers × ad-hoc guardrails",
         "One gateway, one policy schema; a new pipeline is one YAML file and one header",
         "5.0 production and ~211 documented use cases per large enterprise [M-03]"],
        ["CISO / DPO /\nCompliance",
         "Cannot evidence AI oversight to auditors; direct statutory exposure",
         "Append-only audit trail, per-jurisdiction policy packs, HITL decision lineage",
         "₹250 crore DPDP exposure [M-02]; €35M / 7% EU exposure [M-01]"],
        ["Risk / model\ngovernance (BFSI)",
         "Regulator expects demonstrable human-in-the-loop supervision",
         "Escalation queue, fail-closed posture, versioned policy diffs as change evidence",
         "18% of AI systems high-risk, 40% unclassified [M-12]"],
        ["FinOps /\nengineering finance",
         "Inference spend unattributed and unmanaged",
         "Per-pipeline budget policy and tiered routing (enforcement in P1)",
         "15–25× flagship-to-mini price premium [M-06]"],
        ["Head of AI\nprogramme",
         "Pilots stall before production; projects cancelled",
         "Governance as a precondition rather than an afterthought",
         "80%+ pilot failure [M-03]; 40% cancellation without governance [M-12]"],
    ]
    table(s, ML, 2.50, CW, rows, [2.5, 4.6, 5.6, 4.2], pt=11, header_pt=11.5,
          row_h=1.16, header_h=0.50)

    band = rect(s, ML, 8.62, CW, 1.10, GREY_BG, stroke=RULE)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.28)
    p = para(tf, "", 12.5, first=True)
    p.runs[0].text = ""
    style_run(p.add_run(), 12.5, bold=True, color=DEEP).text = "Initial ideal customer profile.  "
    style_run(p.add_run(), 12.5, color=BLACK).text = (
        "Mid-to-large enterprises running three or more production genAI use cases in "
        "regulated or brand-sensitive contexts (BFSI, healthcare, telecommunications), "
        "and India's Global Capability Centers, where 1,600+ centres [M-09] sit at the "
        "intersection of DPDP obligations [M-02] and their parents' EU/US obligations [M-01].")
    source_line(s, "Sources: McKinsey / Gartner [M-03] · DPDP [M-02] · EU AI Act [M-01] "
                   "· appliedAI [M-12] · OpenAI published pricing [M-06] · NASSCOM [M-09]",
                y=SRC_Y + 0.02)


def s_evidence(s):
    slide_title(s, "Every number, with its source", kicker="Evidence")
    lb = textbox(s, ML, 2.40, 9.55, 0.38)
    para(lb.text_frame, "DETECTION QUALITY: INCLUDING WHAT WE MISSED", 11, bold=True,
         color=DEEP, first=True)
    rows = [
        ["Detector", "Blind\n(first contact)", "After disclosed\nrevision", "Target"],
        ["Tier-1 PII recall", "0.836", "0.885", "0.95 (MISSED)"],
        ["Tier-1 PII precision", "1.000", "1.000", "none set"],
        ["Unsourced numeric claims precision", "0.267", "0.857", "none set"],
        ["Prompt injection", "P 1.000 / R 0.150", "not tuned", "none set"],
        ["Toxicity, high band", "P 0.400 / R 0.250", "not tuned", "none set"],
        ["End-to-end verdict accuracy", "n/a", "0.825–0.851\nover 194 cases", "none set"],
    ]
    table(s, ML, 2.86, 9.55, rows, [4.3, 2.0, 2.0, 1.7], pt=11, header_pt=10.5,
          row_h=0.62, header_h=0.66, align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                                            PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    rb = textbox(s, 10.62, 2.40, 8.62, 0.38)
    para(rb.text_frame, "LATENCY: INCLUDING THE PUBLISHED BREACH", 11, bold=True,
         color=DEEP, first=True)
    rows2 = [
        ["Series", "Measured (P99)", "Target", "Verdict"],
        ["Input hold (perceived)", "27.49 ms", "< 50 ms", "Met"],
        ["Sentence hold (perceived)", "39.59 ms", "< 100 ms", "Met"],
        ["tier1_blocklist", "0.020 ms", "2 ms", "Met"],
        ["tier1_pii", "0.133 ms", "2 ms", "Met"],
        ["numeric_claims", "0.277 ms", "5 ms", "Met"],
        ["tier2_toxicity", "23.741 ms", "25 ms", "Met"],
        ["tier2_injection", "25.348 ms", "25 ms", "BREACH (published)"],
    ]
    table(s, 10.62, 2.86, 8.62, rows2, [3.2, 1.6, 1.2, 2.6], pt=11, header_pt=10.5,
          row_h=0.535, header_h=0.66, align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                                             PP_ALIGN.CENTER, PP_ALIGN.LEFT])

    y = 7.60
    lb2 = textbox(s, ML, y, 9.55, 0.38)
    para(lb2.text_frame, "DECISION LAYER AND RELIABILITY", 11, bold=True, color=DEEP,
         first=True)
    tiles = [("840/840", "verdicts correct, 280 cases × 3 policies", "[B-01]"),
             ("7/7", "injected defects caught by the harness", "[B-01]"),
             ("39/39", "fault assertions, 5 quiet-host runs", "[B-07]"),
             ("0", "raw PII in audit evidence", "[B-11]")]
    x = ML
    for big, cap, tag in tiles:
        kpi_tile(s, x, y + 0.44, 2.29, 1.62, big, cap, tag=tag, fill=TINT2, big_pt=19,
                 cap_pt=9.5)
        x += 2.42

    lb3 = textbox(s, 10.62, y, 8.62, 0.38)
    para(lb3.text_frame, "ENGINEERING BASE", 11, bold=True, color=DEEP, first=True)
    tiles2 = [("1,200+", "tests green, two machines and CI", "[B-09]"),
              ("280", "frozen corpus cases, digest-verified", "[B-09]"),
              ("36 / 32 / 9", "ADRs / deviations / limitations", "[B-09]"),
              ("72/72", "figures re-derived from source", "[B-09]")]
    x = 10.62
    for big, cap, tag in tiles2:
        kpi_tile(s, x, y + 0.44, 2.05, 1.62, big, cap, tag=tag, fill=TINT2, big_pt=17,
                 cap_pt=9.5)
        x += 2.19
    source_line(s, REPRO_EVAL + "  ·  " + REPRO_LAT +
                "  ·  python -m eval.fault_injection --reps 5", y=SRC_Y + 0.10)


def s_misses(s):
    slide_title(s, "Why the misses are the point", kicker="Measurement discipline")
    lead(s, "One rule, adopted at the start and not broken under deadline pressure: no "
            "number ships unless a command in the repository reproduces it, and no "
            "unflattering number is removed once measured.")
    cards = [
        ("The PII target was missed and the target did not move.",
         "Recall 0.885 against a 0.95 target. All seven residual misses fall into one "
         "documented scope exclusion. The blind figure of 0.836 ships permanently next "
         "to the revised one. [B-02]"),
        ("Injection recall of 0.150 is weak, and we publish it.",
         "A first-contact number for an untuned model layer sitting behind a pattern "
         "layer that already stops overt attacks. Given prompt injection's OWASP #1 "
         "status [M-05], this is our most important P1 work. [B-04]"),
        ("The end-to-end score fell as coverage rose.",
         "From 0.981 over 159 cases to 0.825–0.851 over 194. The system did not get "
         "worse; the measurement got broader, admitting the hardest classes. Reporting "
         "the improvement without the coverage change would have misled. [B-10]"),
        ("Our calibration procedure inverted on our own data.",
         "The conformal threshold band inverted at every α tried, across 5 seeds. We did "
         "not re-pick α to force a clean result: selecting a parameter after seeing it "
         "fail is tuning toward an outcome. [B-12]"),
    ]
    w, gap = 4.44, 0.24
    x = ML
    for head, body in cards:
        card = rect(s, x, 3.55, w, 3.30, WHITE, stroke=RULE)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.24)
        tf.margin_top = Inches(0.22)
        para(tf, head, 13, bold=True, color=DEEP, first=True, line=1.06)
        para(tf, body, 11.5, color=BLACK, space_before=8, line=1.08)
        x += w + gap

    band = rect(s, ML, 7.28, CW, 1.42, ACCENT)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.34)
    para(tf, "A governance product whose own claims cannot be audited is a contradiction. "
             "The corpus was frozen before any detector was tuned; first measurements "
             "were taken blind.", 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, line=1.08)
    tb = textbox(s, ML, 8.92, CW, 0.72)
    para(tb.text_frame, "Two external comparisons make the figures legible: Microsoft "
                        "Presidio benchmarks at 0.431 span recall on AI4Privacy and "
                        "fine-tuned transformers reach 0.855 [M-08]. Different corpus, "
                        "so a like-for-like head-to-head is a P1 commitment.", 11.5,
         color=BLACK, first=True, line=1.06)
    source_line(s, REPRO_EVAL + " · proposal §2.2, §8.2 and §9", y=SRC_Y + 0.06)


def s_business_model(s):
    slide_title(s, "Open core: the gateway is free, the compliance surface is not",
                kicker="Business model")
    rows = [
        ["Tier", "What is included", "Price posture"],
        ["Open source\n(Apache-2.0)",
         "Gateway, policy engine, detectors, audit schema, evaluation harness: the full "
         "prototype lineage",
         "Free: adoption,\ncommunity, credibility"],
        ["Team",
         "Governance Console (hardened), SSO, retention, alerting",
         "US$500–1,000 per production\npipeline per month  [model]"],
        ["Enterprise",
         "Compliance policy packs (EU AI Act, DPDP, RBI/SEBI, HIPAA), audit export, "
         "air-gapped and local-model deployment, SLAs, calibration on customer traffic",
         "US$40,000–120,000 ACV\n[model]"],
    ]
    table(s, ML, 2.52, CW, rows, [3.0, 10.4, 5.1], pt=12, header_pt=11.5, row_h=1.42,
          header_h=0.50, first_col_bold=True)

    band = rect(s, ML, 7.20, CW, 1.06, GREY_BG, stroke=RULE)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.28)
    p = para(tf, "", 12.5, first=True)
    p.runs[0].text = ""
    style_run(p.add_run(), 12.5, bold=True, color=DEEP).text = (
        "Per-pipeline pricing means revenue tracks the customer's own AI adoption curve, ")
    style_run(p.add_run(), 12.5, color=BLACK).text = (
        "currently averaging 5.0 production use cases and growing 101% year over year [M-03].")

    tb = textbox(s, ML, 8.52, CW, 0.94)
    para(tb.text_frame, "WHY THE FREE TIER IS NOT CANNIBALISATION", 11, bold=True,
         color=ACCENT, first=True)
    para(tb.text_frame, "The paid surfaces are the ones enterprises cannot self-serve: "
                        "compliance policy packs, audit retention, per-customer "
                        "calibration, SLAs. The gateway being free is what drives the "
                        "adoption we monetise.", 12, color=BLACK, space_before=4,
         line=1.06)
    source_line(s, "Source: proposal §10.1 and §9 risk 9 · pricing is [model]; adoption "
                   "growth from McKinsey / Gartner [M-03]", y=SRC_Y + 0.02)


def s_unit_economics(s):
    slide_title(s, "Unit economics and a three-year outline", kicker="Business case")
    lb = textbox(s, ML, 2.42, 10.20, 0.38)
    para(lb.text_frame, "UNIT ECONOMICS  [model]", 11, bold=True, color=DEEP, first=True)
    rows = [
        ["Metric", "Assumption", "Basis"],
        ["Blended ACV", "US$35,000–50,000",
         "One enterprise logo ≈ 3–5 paid pipelines plus pack/SLA uplift"],
        ["Gross margin", "~85%",
         "Software plus thin sampling; the customer pays their own model spend"],
        ["CAC (partner-led)", "US$10,000–18,000",
         "SI-attached deals and OSS inbound; no outbound sales team in years 1–2"],
        ["Payback", "< 6 months", "ACV ÷ CAC at the above"],
        ["Expansion", "1 pipeline → 4–6 in 12–18 months",
         "Tracks documented enterprise use-case growth [M-03]"],
    ]
    table(s, ML, 2.88, 10.20, rows, [2.4, 3.4, 4.4], pt=11, header_pt=11, row_h=0.80,
          header_h=0.48, first_col_bold=True)

    rb = textbox(s, 11.32, 2.42, 7.92, 0.38)
    para(rb.text_frame, "THREE-YEAR OUTLINE  [model]", 11, bold=True, color=DEEP,
         first=True)
    years = [
        ("Year 1", "US$150–250k ARR", "4–6 pilot conversions", TINT1, DEEP),
        ("Year 2", "US$1–1.5M ARR", "Partner-led repeatability", ACCENT, WHITE),
        ("Year 3", "US$5–8M ARR", "120–180 accounts at blended ACV", DEEP, WHITE),
    ]
    y = 2.88
    for label, arr, note, fill, fg in years:
        sp = rect(s, 11.32, y, 7.92, 1.28, fill,
                  stroke=RULE if fill is TINT1 else None)
        tf = sp.text_frame
        tf.margin_left = tf.margin_right = Inches(0.26)
        p = para(tf, "", 12, first=True)
        p.runs[0].text = ""
        style_run(p.add_run(), 11.5, bold=True, color=fg).text = label + "     "
        style_run(p.add_run(), 20, bold=True, color=fg).text = arr
        para(tf, note, 11.5, color=fg, space_before=3)
        y += 1.40

    note = rect(s, 11.32, 7.10, 7.92, 1.52, GREY_BG, stroke=RULE)
    tf = note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = Inches(0.18)
    para(tf, "A DELIBERATELY MODEST SHARE ASSUMPTION", 10.5, bold=True, color=ACCENT,
         first=True)
    para(tf, "The Year-3 figure is between 0.35% and 0.56% of the narrower analyst "
             "estimate for 2030 [M-04].", 12, color=BLACK, space_before=6, line=1.06)

    cost = rect(s, ML, 7.10, 10.20, 1.86, DEEP)
    tf = cost.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = Inches(0.18)
    para(tf, "ON COST SAVINGS, WE PUBLISH NO ABSOLUTE FIGURE", 10.5, bold=True,
         color=ACCENT, first=True)
    para(tf, "Published pricing shows a 15–25× flagship-to-mini premium [M-06]; our own "
             "demo configuration runs a 2.0× tier gap, so savings are reported as "
             "relative and scale with the deployment's own ratio [B-13]. Cost-plane "
             "enforcement is P1: claiming a dollar number we have not measured would "
             "violate the principle this deck is built on.", 11.5, color=WHITE,
         space_before=6, line=1.06)
    source_line(s, "Source: proposal §10.2–10.3 · [model] figures are assumptions, "
                   "labelled as such · tier pricing in config/gateway.yaml [B-13]")


def s_gtm(s):
    slide_title(s, "Open source in, systems integrators out", kicker="Go-to-market")
    phases = [
        ("LAUNCH", "0–3 months",
         "Open-source release plus the build-in-public story: the adjudication ledger "
         "and blind-measurement discipline are genuinely distinctive developer content; "
         "conference submissions on “AI built under AI oversight”.",
         "Top-of-funnel attention; inbound design-partner pipeline", TINT1, DEEP),
        ("DESIGN PARTNERS", "2–6 months",
         "3–5 pilots across Indian BFSI and GCCs [M-09] plus one EU-exposed SaaS [M-01]. "
         "Success is one pipeline in production with audit evidence accepted by the "
         "partner's risk function.",
         "Reference architecture; first case studies; calibration playbook", ACCENT, WHITE),
        ("PARTNER-LED", "6–18 months",
         "Systems-integrator channel. Accenture alone has committed US$3 billion over "
         "three years to Data & AI and doubled its practice to 80,000 specialists with "
         "Responsible AI frameworks across 19 industries [M-10].",
         "Repeatable SI-attached deals; first US$1M ARR [model]", DEEP, WHITE),
    ]
    w, gap = 5.92, 0.36
    x = ML
    for label, when, body, proof, fill, fg in phases:
        card = rect(s, x, 2.50, w, 4.30, fill, stroke=RULE if fill is TINT1 else None)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.28)
        tf.margin_top = Inches(0.24)
        para(tf, label, 17, bold=True, color=fg, first=True)
        para(tf, when, 11.5, bold=True, color=fg, space_before=2)
        para(tf, body, 12, color=fg, space_before=10, line=1.08)
        para(tf, "PROOF POINT", 10, bold=True, color=fg, space_before=12)
        para(tf, proof, 11.5, color=fg, space_before=2, line=1.06)
        x += w + gap

    band = rect(s, ML, 7.16, CW, 1.34, GREY_BG, stroke=RULE)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.30)
    para(tf, "The flywheel", 12, bold=True, color=ACCENT, first=True)
    para(tf, "Open-source adoption produces policy patterns and calibration data  →  "
             "patterns become packaged compliance packs  →  packs make SI deployments "
             "faster  →  SI deployments produce references that pull further adoption.",
         12.5, color=BLACK, space_before=4, line=1.06)
    source_line(s, "Sources: proposal §10.5 · NASSCOM GCC reporting [M-09] · EU AI Act "
                   "[M-01] · Accenture Newsroom, programme confirmed active in 2025 "
                   "releases [M-10]", y=SRC_Y - 0.02)


def s_impact(s):
    slide_title(s, "What changes for the enterprise buyer", kicker="Impact")
    rows = [
        ["Benefit", "Mechanism", "Supporting evidence"],
        ["Prevent incidents rather than discover them",
         "Verdicts are issued before delivery; a flagged sentence never partially reaches "
         "the user",
         "840/840 conformance [B-01]; sentence-buffered interception [B-06]"],
        ["The provider never sees regulated data",
         "Input-lane PII redaction happens pre-dispatch",
         "Live behaviour; 0 raw values in audit evidence [B-11]"],
        ["Evidence for the regulator, as a by-product",
         "Append-only records, categories not values, HITL lineage, policy version stamped "
         "per record",
         "Direct response to AI Act oversight duties [M-01] and DPDP safeguards [M-02]"],
        ["Governance without redeployment",
         "Policies are versioned YAML, hot-reloadable, per pipeline and per jurisdiction",
         "Three pipelines diverge with zero code difference [B-01]"],
        ["Failure has a defined posture",
         "Fail-open versus fail-closed is a policy field; fail-closed escalates to a human "
         "and never silently blocks",
         "39/39 fault assertions across 5 runs [B-07]"],
        ["No cloud lock-in",
         "Provider-neutral, including fully local and air-gapped operation",
         "Demonstrated against both a cloud provider and a local model"],
    ]
    table(s, ML, 2.50, CW, rows, [4.3, 8.0, 6.2], pt=11.5, header_pt=11.5, row_h=1.00,
          header_h=0.50)

    band = rect(s, ML, 8.72, CW, 1.00, DEEP)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.32)
    para(tf, "Three clocks are running at once: regulatory (fines already active), "
             "market (35–45% CAGR), operational (80%+ of pilots stalled). The window "
             "for a converged governance layer is open now.", 13.5, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, first=True, line=1.06)
    source_line(s, "Source: proposal §14 · " + REPRO_EVAL, y=SRC_Y + 0.02)


def s_roadmap(s):
    slide_title(s, "P0 is done. P1 closes the gaps we measured.", kicker="Roadmap")
    D.roadmap_timeline(s)
    source_line(s, "Source: proposal §11 · exit criteria: P1 recall ≥ 0.95 on the "
                   "expanded corpus or a documented explanation, 3 pilots live · "
                   "P2 first paid conversions · P3 partner-led revenue")


def s_risks(s):
    slide_title(s, "Risks, led by our own weakest numbers", kicker="Risk analysis")
    rows = [
        ["#", "Risk", "Honest reality", "Mitigation"],
        ["1", "Detection quality is incomplete",
         "Injection recall 0.150 blind; toxicity high-band recall 0.250 [B-04]",
         "Deterministic patterns already stop overt attacks at microsecond cost [B-05]. "
         "Thresholds are designed to be tuned on real traffic, which is the paid enterprise motion."],
        ["2", "Calibration failed on our own data",
         "The conformal threshold procedure inverted at every α tried, across 5 seeds [B-12]",
         "The harness worked: it detected that the grounding signal does not yet separate "
         "classes on a synthetic corpus. Real-traffic calibration in P1."],
        ["3", "Hyperscaler bundling",
         "AWS, Azure and NVIDIA ship guardrail engines with per-use-case policies and "
         "fail posture [M-07]",
         "Compete on convergence, neutrality and evidence, not on features they already "
         "have. Cross-provider and air-gapped operation is structurally unavailable to them."],
        ["4", "Latency objections",
         "One detector breaches its budget at P99 [B-05]; full-coverage windowing costs "
         "more on long inputs",
         "Deterministic tier is three orders of magnitude inside budget [B-05]; "
         "user-perceived holds meet published expectations [M-11]. Compiled data plane in P2."],
        ["5", "Synthetic-corpus overfit critique",
         "Fair for any self-built evaluation. Our corpus is ours",
         "Frozen before tuning; blind-first protocol; independently reviewed. P1 commits "
         "to a Presidio head-to-head [M-08] and calibration on partner traffic."],
        ["6", "Provider churn",
         "A model tier we depended on was deprecated 11 days before this submission",
         "It failed loudly, not silently. Provenance-classed providers and the accounting "
         "canary are shipped features that exist because of this class of event."],
        ["7", "Market timing / category risk",
         "Analyst estimates diverge 4.08× on market size [M-04]; consolidation is under "
         "way [M-07]",
         "Consolidation validates the category. Open-core distribution reduces dependence "
         "on any single forecast; both estimates agree on 35%+ CAGR."],
        ["8", "Team capacity",
         "Three third-year undergraduates, no enterprise delivery history",
         "A complete public delivery record most professional teams cannot produce [B-09]. "
         "The SI channel [M-10] supplies deployment scale."],
    ]
    table(s, ML, 2.44, CW, rows, [0.5, 3.2, 5.3, 9.5], pt=10.5, header_pt=11,
          row_h=0.86, header_h=0.44)
    source_line(s, "Source: proposal §9 (nine risks; the ninth, open-core "
                   "cannibalisation, is addressed on the business-model slide) · "
                   + REPRO_EVAL, y=SRC_Y + 0.10)


def s_team(s):
    slide_title(s, "Why this team: the delivery record", kicker="Qualifications")
    lead(s, "We have no enterprise track record to present. In place of one, we offer a "
            "complete public engineering record: a reviewer can inspect our judgement "
            "directly rather than infer it from a résumé.")
    tiles = [("36", "architecture decision records", "[B-09]"),
             ("32", "deviations filed and all closed", "[B-09]"),
             ("61", "minor resolutions logged", "[B-09]"),
             ("9", "standing limitations published", "[B-09]"),
             ("4", "times the AI reviewer blocked the build", "[B-09]"),
             ("72/72", "figures re-derived from source artifacts", "[B-09]")]
    w, gap = 2.95, 0.13
    x = ML
    for big, cap, tag in tiles:
        kpi_tile(s, x, 3.36, w, 1.86, big, cap, tag=tag, fill=TINT2, big_pt=23,
                 cap_pt=10)
        x += w + gap

    lb = textbox(s, ML, 5.62, CW, 0.38)
    para(lb.text_frame, "WHAT WE DID WHEN THINGS WENT WRONG: ON THE RECORD", 11,
         bold=True, color=DEEP, first=True)
    rows = [
        ["Incident", "Resolution", "Product consequence"],
        ["A third-party gateway silently injected ~5,000 hidden prompt tokens per "
         "request, corrupting every cost measurement",
         "Detected by a boot-time canary comparing provider accounting against a local "
         "estimate [B-14]",
         "Provenance-classed providers and the accounting canary became shipped features"],
        ["Wall-clock budget enforcement manufactured spurious detector faults under "
         "load, enough to flip a verdict in a fail-closed pipeline",
         "Rebuilt enforcement around detector-attributable CPU time; spurious faults "
         "2/30 → 0/30 [B-08]",
         "Attribution-correct enforcement, now a differentiator"],
        ["Market research contradicted our own differentiation claim",
         "Claim retracted, positioning rebuilt on ground that survives the evidence [M-07]",
         "A more defensible five-point differentiation"],
    ]
    table(s, ML, 6.08, CW, rows, [6.2, 6.2, 6.1], pt=11, header_pt=11, row_h=1.06,
          header_h=0.46)

    band = rect(s, ML, 9.36, CW, 0.70, ACCENT)
    tf = band.text_frame
    tf.margin_left = tf.margin_right = Inches(0.30)
    para(tf, "For a company whose product is trust in AI systems, a timestamped record "
             "of how we behave when we are wrong is the qualification.", 13, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, first=True)
    source_line(s, "Source: proposal §12–13 · docs/03-decisions.md, "
                   "docs/08-open-questions.md · python -m eval.check_derivations",
                y=SRC_Y + 0.28)


def s_ask(s):
    slide_title(s, "The ask: INR 1.6 crore for 12 months", kicker="The ask")
    lead(s, "Indicative pre-seed of INR 1.6 crore (~US$190,000), to convert this "
            "prototype into three design-partner pilots within two quarters. [model]")
    D.use_of_funds(s)

    lb = textbox(s, ML, 6.30, CW, 0.38)
    para(lb.text_frame, "MILESTONE GATES", 11, bold=True, color=DEEP, first=True)
    rows = [
        ["Quarter", "Milestone gate"],
        ["Q1", "Open-source launch; 3 design-partner pilots signed; Console GA track scoped"],
        ["Q2", "3 pilots live; first partner-traffic calibration cycles; expanded-corpus "
               "recall target met or explained"],
        ["Q3", "Threshold and quality targets met on partner traffic; first paid conversions"],
        ["Q4", "SOC 2 kickoff; repeatable SI-attached deal template; year-2 plan grounded "
               "in pilot data"],
    ]
    table(s, ML, 6.76, CW, rows, [1.5, 16.98], pt=12, header_pt=11.5, row_h=0.58,
          header_h=0.46, first_col_bold=True)
    source_line(s, "Source: proposal §10.4 · the funding figure and allocation are "
                   "[model] assumptions, stated as such", y=SRC_Y + 0.10)


def s_closing(s):
    rect(s, 0, 0, 20.0, 11.25, DEEP)
    rect(s, 0, 0, 0.34, 11.25, ACCENT)
    tb = textbox(s, 2.10, 1.85, 15.6, 1.90)
    para(tb.text_frame, "We built the supervision layer,", 34, bold=True, color=WHITE,
         first=True, line=1.02)
    para(tb.text_frame, "and it runs.", 34, bold=True, color=ACCENT, line=1.02)

    items = [
        ("840/840", "verdicts correct on a frozen corpus [B-01]"),
        ("0.133 ms", "deterministic detection at P99 [B-05]"),
        ("1,200+", "tests on two machines and CI [B-09]"),
        ("9", "published limitations, documented rather than hidden"),
    ]
    y = 4.30
    for big, cap in items:
        bb = textbox(s, 2.10, y, 3.05, 0.52)
        para(bb.text_frame, big, 20, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT,
             first=True)
        cb = textbox(s, 5.40, y + 0.06, 8.60, 0.46)
        para(cb.text_frame, cap, 13, color=WHITE, first=True)
        y += 0.66

    lb = textbox(s, 2.10, 7.30, 6.60, 0.36)
    para(lb.text_frame, "WHAT WE ASK OF THIS PANEL", 11, bold=True, color=TINT1,
         first=True)
    asks = [
        "Select ControlPlane to advance in the Accenture Innovation Challenge 2026.",
        "Open the design-partner conversation inside Accenture's responsible-AI practice "
        "backed by US$3 billion and 80,000 specialists [M-10].",
        "Consider the indicative pre-seed of INR 1.6 crore to convert this prototype "
        "into three live pilots within two quarters.",
    ]
    y = 7.78
    for i, a in enumerate(asks, 1):
        nb = textbox(s, 2.10, y, 0.50, 0.44)
        para(nb.text_frame, str(i), 15, bold=True, color=ACCENT, first=True)
        tb2 = textbox(s, 2.66, y, 11.30, 0.72)
        para(tb2.text_frame, a, 12.5, color=WHITE, first=True, line=1.06)
        y += 0.80

    card = rect(s, 14.35, 4.30, 3.75, 3.05, ACCENT)
    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = Inches(0.24)
    para(tf, "THE TEST WE INVITE", 10.5, bold=True, color=WHITE, first=True)
    para(tf, "Clone the repository, run one command, and check any number in this deck.",
         14, bold=True, color=WHITE, space_before=8, line=1.08)
    para(tf, "python -m eval.run_all", 11.5, color=WHITE, space_before=10)
    para(tf, "We built it so that you could.", 11.5, italic=True, color=WHITE,
         space_before=10)

    fb = textbox(s, 2.10, 10.05, 15.6, 0.40)
    para(fb.text_frame, "ControlPlane: oversight for AI, built under oversight.", 15,
         bold=True, color=WHITE, first=True)


# ===================================================================== assembly
DIVIDERS = {
    "problem": ("01", "The problem",
                "Three failures, one root cause, and a regulatory clock that has already started."),
    "solution": ("02", "The solution",
                 "One URL change, three planes of detection, and exactly one verdict per response."),
    "market": ("03", "Market and position",
               "Two analyst estimates, an empty quadrant, and one claim we retracted."),
    "business": ("04", "Business and impact",
                 "Open core, partner-led distribution, and a roadmap that is our limitations register."),
    "ask": ("05", "The ask",
            "What we want, what it buys, and how you can check every number in this deck."),
}

# (kind, payload) in final presentation order, between the team slide and Thank you
FLOW = [
    ("content", s_exec_summary),
    ("divider", "problem"),
    ("content", s_problem),
    ("content", s_problem_evidence),
    ("content", s_regulatory),
    ("divider", "solution"),
    ("content", s_solution_overview),
    ("content", s_architecture),
    ("content", s_one_signal),
    ("content", s_four_verdicts),
    ("content", s_signature),
    ("content", s_escalation),
    ("content", s_ships),
    ("divider", "market"),
    ("content", s_market),
    ("content", s_competitive),
    ("content", s_retraction),
    ("content", s_users),
    ("content", s_evidence),
    ("content", s_misses),
    ("divider", "business"),
    ("content", s_business_model),
    ("content", s_unit_economics),
    ("content", s_gtm),
    ("content", s_impact),
    ("content", s_roadmap),
    ("content", s_risks),
    ("content", s_team),
    ("divider", "ask"),
    ("content", s_ask),
    ("content", s_closing),
]

DARK_SLIDES = {s_closing}


def main():
    prs = Presentation(str(TEMPLATE))
    deck = Deck(prs, prs.slide_layouts[6])

    title_slide, team_slide, thankyou_slide = prs.slides[0], prs.slides[2], prs.slides[6]
    fix_title_slide(title_slide)
    fill_team_slide(team_slide)

    n_template_before = len(prs.slides._sldIdLst)  # 7
    number = 3  # title = 1, team = 2
    for kind, payload in FLOW:
        s = deck.new()
        if kind == "divider":
            num, title, sub = DIVIDERS[payload]
            divider(s, num, title, sub)
            footer(s, number, dark=True)
        else:
            payload(s)
            footer(s, number, dark=payload in DARK_SLIDES)
        number += 1

    fix_thankyou_slide(thankyou_slide, number)

    bold_italic_from_facename(prs)
    force_arial(prs)

    content_idx = list(range(n_template_before, n_template_before + len(FLOW)))
    reorder_and_prune(prs, [0, 2] + content_idx + [6])
    prs.save(str(OUT))

    check = Presentation(str(OUT))
    print(f"saved {OUT.name}: {len(check.slides)} slides, "
          f"{OUT.stat().st_size / 1_048_576:.2f} MB")
    return len(check.slides)


if __name__ == "__main__":
    main()
