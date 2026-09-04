"""Merge the staged content slides into the locked AIC template package.

Direction matters: the TEMPLATE is the base and is never rewritten. The cover and
thank-you slides are picture-fill freeforms whose fidelity depends on the template's
own masters, theme and media parts, so dragging them into a generated package is how
"byte-identical in layout" breaks. Generated content slides carry no media except one
screenshot, which is re-added through the normal picture API rather than by XML.
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).resolve().parent))
from deck_lib import bold_italic_from_facename  # noqa: E402
# The locked shells are filled by the SAME functions the previous build used, so the
# cover/team/thank-you text has one source rather than two diverging copies.
from build_deck import (fix_title_slide, fill_team_slide,  # noqa: E402
                        fix_thankyou_slide)

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "assets" / "AIC_Talent-Brand_PPT-Template.pptx"
STAGED = Path("/tmp/cap/deck_v2_content.pptx")
OUT = ROOT / "submission" / "b24bb1029_ControlPlane.pptx"

INK = "17171C"
BODY = "212121"
MUTED = "93939F"
SLATE = "75758A"
ACCENT = "A100FF"

# 200 words max each, per the template's own instruction. Counted in the assert below.
PROBLEM_200 = (
    "Enterprises have adopted generative AI faster than they can supervise it. Seventy-two "
    "percent of organisations now use it in at least one function, yet more than eighty percent "
    "of pilots fail to reach production and forty percent of initiatives are cancelled where "
    "governance is absent.\n"
    "The failures are not exotic. A support assistant states a refund that never happened. An "
    "HR copilot repeats an employee's phone number into a retained transcript. An advisory tool "
    "loops on a malformed integration, burning inference budget for a week. Three planes, "
    "quality, cost and responsibility, and three teams discovering three incidents weeks later "
    "from three different systems.\n"
    "The common cause is structural. Observability platforms report after delivery. Guardrail "
    "engines filter one category of harm. Gateways route traffic. None makes a single "
    "accountable decision about a response before delivery, and none lets one organisation hold "
    "a customer chatbot and an internal HR tool to different standards without shipping "
    "different code.\n"
    "Meanwhile the regulatory clock runs. The EU AI Act reaches thirty-five million euros or "
    "seven percent of global turnover, India's DPDP Act two hundred and fifty crore rupees, and "
    "the obligation lands on the deploying enterprise, not the model vendor."
)

SOLUTION_200 = (
    "ControlPlane is a reverse-proxy oversight gateway. An application changes one base URL, "
    "keeping its own SDK. Every request and response then passes detection lanes scoring three "
    "planes, converging to one verdict before delivery: pass, edit, block or escalate.\n"
    "Interception is per sentence rather than per token or per finished response, so a flagged "
    "sentence never partially reaches the user while the hold stays inside a measured budget. "
    "Deterministic matchers run three orders of magnitude inside their two millisecond budget. "
    "Quantized CPU classifiers handle injection and toxicity. Anything expensive, semantic "
    "entropy, fairness auditing, model-judge scoring, runs on an asynchronous lane and is never "
    "awaited on the hot path.\n"
    "The decisive property is that policy lives in versioned YAML, one file per use case, rather "
    "than in code. The same sentence carrying the same three detector labels resolves to edit "
    "under a support policy, block under an HR policy and escalate under a financial one. That "
    "is verified by running the engine, and a test fails if any use-case name reaches executable "
    "code.\n"
    "Escalated responses are quarantined rather than delivered, and every decision is written "
    "into an append-only audit lineage recording the category of an interception, never the "
    "matched value."
)


MONO = "Courier New"
TEXT = "Arial"


def force_deck_fonts(prs):
    """Normalise typefaces to the two the design system allows, and no others.

    `deck_lib.force_arial` rewrites *every* face to Arial. That is right for the
    Arial-only deck it was written for, and wrong here: this deck reserves Courier New
    for identifiers, mono kickers, evidence tags and commands, so a blanket sweep
    silently erases half the type system. Run-level Courier is preserved; everything
    else, including theme fonts, `+mj-lt`/`+mn-lt` references and the template's Arimo
    `a:sym`, still collapses to Arial so no third face can survive in the package.
    """
    from pptx.oxml.ns import qn

    docs = [prs.part._element]
    for m in prs.slide_masters:
        docs.append(m._element)
        for lay in m.slide_layouts:
            docs.append(lay._element)
    for s in prs.slides:
        docs.append(s._element)
        if s.has_notes_slide:
            docs.append(s.notes_slide._element)
    theme_docs = []
    for part in prs.part.package.iter_parts():
        name = str(part.partname)
        if "/theme/" in name or "notesMaster" in name:
            el = getattr(part, "_element", None)
            if el is not None:
                theme_docs.append(el)

    tags = (qn("a:latin"), qn("a:ea"), qn("a:cs"), qn("a:sym"))

    # Slides/masters/layouts: keep Courier where it was deliberately set.
    for root in docs:
        for tag in tags:
            for el in root.iter(tag):
                face = el.get("typeface") or ""
                el.set("typeface", MONO if face == MONO else TEXT)
                for junk in ("panose", "pitchFamily", "charset"):
                    el.attrib.pop(junk, None)
        for el in list(root.iter()):
            v = el.get("typeface")
            if v and v.startswith("+"):
                el.set("typeface", TEXT)

    # Theme parts carry no run-level intent: everything becomes Arial.
    for root in theme_docs:
        for tag in tags:
            for el in root.iter(tag):
                el.set("typeface", TEXT)
                for junk in ("panose", "pitchFamily", "charset"):
                    el.attrib.pop(junk, None)
        for el in list(root.iter()):
            v = el.get("typeface")
            if v and v.startswith("+"):
                el.set("typeface", TEXT)


def _wc(text: str) -> int:
    return len(text.replace("\n", " ").split())


def clear_and_write(slide, text, *, x, y, w, h, pt, color=BODY, line_mult=1.16):
    """Add a body textbox. The template's own heading group is left untouched."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    from pptx.util import Emu as _E
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = text.split("\n")
    for i, para_text in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(pt)
        run.font.name = "Arial"
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor.from_string(color)
        p.line_spacing = line_mult
        if i < len(paras) - 1:
            p.space_after = Pt(9)
    return box


def add_page_number(slide, n, *, dark=False):
    from pptx.dml.color import RGBColor
    box = slide.shapes.add_textbox(Inches(18.9), Inches(10.7), Inches(0.35), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(10)
    r.font.name = "Arial"
    r.font.color.rgb = RGBColor.from_string("B9B9C2" if dark else MUTED)


def add_footer_note(slide, text, *, mono=False, dark=False):
    from pptx.dml.color import RGBColor
    box = slide.shapes.add_textbox(Inches(0.75), Inches(10.3), Inches(17), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.name = "Courier New" if mono else "Arial"
    r.font.color.rgb = RGBColor.from_string("B9B9C2" if dark else MUTED)


def copy_slide_shapes(src_slide, dst_slide, *, pictures):
    """Copy every shape by XML, except pictures which are re-added via the API."""
    for shape in src_slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            key = (round(Emu(shape.left).inches, 2), round(Emu(shape.top).inches, 2))
            path = pictures.get(key)
            if path is None:
                raise SystemExit(f"unmapped picture at {key} - refusing to drop it")
            dst_slide.shapes.add_picture(
                str(path), shape.left, shape.top, shape.width, shape.height
            )
            continue
        dst_slide.shapes._spTree.append(copy.deepcopy(shape._element))


P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def set_background(dst_slide, src_slide):
    """Carry a solid slide background across (dividers, product band, close).

    `p:bg` is a child of `p:cSld`, never of `p:sld`. Looking for it one level too high
    returns None for every slide, so this function returned early every time and silently
    dropped all 24 backgrounds: the six dark slides shipped as white text on white. Schema
    order also matters on the way in, `p:bg` must precede `p:spTree`, hence insert(0).
    """
    src_cSld = src_slide._element.find(f"{P_NS}cSld")
    if src_cSld is None:
        raise SystemExit("source slide has no p:cSld")
    src_bg = src_cSld.find(f"{P_NS}bg")
    if src_bg is None:
        return False
    dst_cSld = dst_slide._element.find(f"{P_NS}cSld")
    if dst_cSld is None:
        raise SystemExit("destination slide has no p:cSld")
    existing = dst_cSld.find(f"{P_NS}bg")
    if existing is not None:
        dst_cSld.remove(existing)
    dst_cSld.insert(0, copy.deepcopy(src_bg))
    return True


def main():
    assert _wc(PROBLEM_200) <= 200, f"problem statement is {_wc(PROBLEM_200)} words"
    assert _wc(SOLUTION_200) <= 200, f"solution statement is {_wc(SOLUTION_200)} words"
    print(f"word counts: problem={_wc(PROBLEM_200)}  solution={_wc(SOLUTION_200)}")

    tpl = Presentation(str(TEMPLATE))
    stg = Presentation(str(STAGED))
    blank = tpl.slide_layouts[6]

    pictures = {
        (8.75, 2.55): ROOT / "reports/deck-render/console-capture-dashboard.png",
    }

    n_tpl = len(tpl.slides._sldIdLst)
    assert n_tpl == 7, f"template has {n_tpl} slides, expected 7"

    # Copy each staged content slide into a new blank slide.
    bg_moved = 0
    for src in stg.slides:
        dst = tpl.slides.add_slide(blank)
        bg_moved += bool(set_background(dst, src))
        copy_slide_shapes(src, dst, pictures=pictures)
    print(f"copied {len(stg.slides)} content slides, {bg_moved} backgrounds")
    if bg_moved != len(stg.slides):
        raise SystemExit(
            f"only {bg_moved}/{len(stg.slides)} backgrounds transferred; "
            "white-on-white slides would ship"
        )

    ids = list(tpl.slides._sldIdLst)
    tpl_cover, tpl_instr, tpl_team = ids[0], ids[1], ids[2]
    tpl_prob, tpl_soln, tpl_video, tpl_thanks = ids[3], ids[4], ids[5], ids[6]
    content = ids[7:]
    assert len(content) == 24, f"expected 24 content slides, got {len(content)}"

    # Final order per docs/deck-plan.md §1.
    order = [
        tpl_cover,        # 1  cover
        tpl_team,         # 2  team details
        content[0],       # 3  executive summary
        content[1],       # 4  divider 01
        tpl_prob,         # 5  problem statement (200 words)
        content[2],       # 6  three failures
        content[3],       # 7  regulatory clock
        content[4],       # 8  divider 02
        tpl_soln,         # 9  proposed solution (200 words)
        content[5],       # 10 request lifecycle
        content[6],       # 11 one signal
        content[7],       # 12 signature
        content[8],       # 13 console
        content[9],       # 14 escalation ladder
        content[10],      # 15 divider 03
        content[11],      # 16 what ships
        content[12],      # 17 detection quality
        content[13],      # 18 latency and failure
        content[14],      # 19 why the misses
        content[15],      # 20 divider 04
        content[16],      # 21 two analyst estimates
        content[17],      # 22 competitive
        content[18],      # 23 who signs
        content[19],      # 24 open core
        content[20],      # 25 gtm and roadmap
        content[21],      # 26 risks
        content[22],      # 27 references
        content[23],      # 28 the ask
        tpl_video,        # 29 video
        tpl_thanks,       # 30 thank you
    ]
    assert len(order) == 30, f"order has {len(order)} entries"

    sldIdLst = tpl.slides._sldIdLst
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in order:
        sldIdLst.append(el)
    # The instruction slide is simply never re-appended: dropped, per the brief.
    print("instruction slide dropped; 30 slides ordered")

    slides = list(tpl.slides)

    # Locked shells: only the text we are permitted to change.
    fix_title_slide(slides[0])
    fill_team_slide(slides[1])
    fix_thankyou_slide(slides[29], 30)

    # Mandated named slides: keep the template heading, add the body copy.
    # 18.5in at 16pt measured ~142 characters per line and bottomed out at y=4.74,
    # leaving two-thirds of the slide empty (V7 under-fill). The design system caps
    # prose measure at 12in (DESIGN-controlplane-deck.md lede.width) and body-large
    # is 20pt, which measures ~85 chars/line and ends at y=7.26 on the longer of the
    # two blocks. Type grows rather than shrinks, so V7 is satisfied by layout.
    clear_and_write(slides[4], PROBLEM_200, x=0.75, y=2.05, w=12.0, h=7.6, pt=20)
    add_footer_note(slides[4], "Source: McKinsey State of AI 2025; Gartner; "
                               "Regulation (EU) 2024/1689; MeitY DPDP Rules 2025  ·  "
                               "[M-01] [M-02] [M-03] [M-12]")
    add_page_number(slides[4], 5)

    clear_and_write(slides[8], SOLUTION_200, x=0.75, y=2.05, w=12.0, h=7.6, pt=20)
    add_footer_note(slides[8], "Source: docs/02-architecture.md, docs/04-policy-and-detection-spec.md, "
                               "reports/latency_report.md  ·  [B-01] [B-05] [B-06]")
    add_page_number(slides[8], 9)

    add_page_number(slides[28], 29)

    bold_italic_from_facename(tpl)
    force_deck_fonts(tpl)

    tpl.save(str(OUT))
    chk = Presentation(str(OUT))
    print(f"saved {OUT.name}: {len(chk.slides)} slides, "
          f"{OUT.stat().st_size / 1_048_576:.2f} MB")


if __name__ == "__main__":
    main()
