"""Native-shape diagram builders for the ControlPlane deck.

Every diagram here is drawn with python-pptx autoshapes and connectors -- no images.
Geometry is expressed in template inches (20 x 11.25 canvas); font sizes are
brief-relative and scaled by deck_lib.fs().
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from deck_lib import (ACCENT, BLACK, CW, DEEP, GREY_BG, MUTED, ML, RULE, TINT1, TINT2,
                      WHITE, arrow, edge_label, hline, node, para, rect, text_width,
                      textbox)


def _row(x0, x1, count, gap):
    """Equal-width nodes across [x0, x1] with uniform gaps -> list of (x, w)."""
    w = (x1 - x0 - gap * (count - 1)) / count
    return [(x0 + i * (w + gap), w) for i in range(count)]


def _fan_label(slide, sx, sy, ex, ey, ly, lines, *, side, w=3.0, gap=0.16, pt=10):
    """Label a fan connector without letting the line strike through the text.

    A fixed sideways offset fails on steep connectors: the line moves further
    horizontally across the label band than the offset clears. So solve for the
    x-range the line actually occupies over the label's own vertical band, then
    park the box just outside it and hug the text to that edge.
    """
    h = 0.20 * len(lines) + 0.10
    span = (ey - sy) or 1e-6
    xs = [sx + (ex - sx) * ((y - sy) / span) for y in (ly - h / 2, ly + h / 2)]
    if side < 0:
        right = min(xs) - gap
        return edge_label(slide, right - w / 2, ly, "", lines=lines, w=w, pt=pt,
                          align=PP_ALIGN.RIGHT)
    left = max(xs) + gap
    return edge_label(slide, left + w / 2, ly, "", lines=lines, w=w, pt=pt,
                      align=PP_ALIGN.LEFT)


# ------------------------------------------------------------- DIAGRAM 1 · lifecycle
def lifecycle(slide):
    y, h = 3.16, 2.80
    cells = _row(ML, ML + CW, 6, 0.34)
    spec = [
        ("APP", ["One URL change"], TINT1, DEEP),
        ("INPUT LANE", ["PII scrub", "Injection screen"], TINT1, DEEP),
        ("MODEL", ["Any provider,", "tiered"], GREY_BG, BLACK),
        ("SENTENCE BUFFER", ["Checks before", "release"], TINT1, DEEP),
        ("POLICY ENGINE", ["One verdict"], ACCENT, WHITE),
        ("AUDIT + HITL", ["Evidence", "Review"], DEEP, WHITE),
    ]
    for (x, w), (title, body, fill, fg) in zip(cells, spec):
        node(slide, x, y, w, h, title, body, fill=fill, fg=fg, title_pt=13,
             body_pt=11, stroke=RULE if fill in (TINT1, GREY_BG) else None)
    for i in range(5):
        x_end = cells[i][0] + cells[i][1]
        arrow(slide, x_end + 0.05, y + h / 2, cells[i + 1][0] - 0.05, y + h / 2)

    pb = textbox(slide, ML, y + h + 0.72, CW, 0.32)
    para(pb.text_frame, "FOUR VERDICTS. ONE ACCOUNTABLE DECISION PER RESPONSE.", 10.5,
         bold=True, color=MUTED, first=True)

    py, ph = y + h + 1.28, 1.24
    pills = _row(ML, ML + CW, 4, 0.40)
    pspec = [("PASS", TINT1, DEEP), ("EDIT", ACCENT, WHITE),
             ("BLOCK", DEEP, WHITE), ("ESCALATE", ACCENT, WHITE)]
    for (x, w), (label, fill, fg) in zip(pills, pspec):
        sp = rect(slide, x, py, w, ph, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                  radius=0.5, stroke=RULE if fill is TINT1 else None)
        para(sp.text_frame, label, 15, bold=True, color=fg, align=PP_ALIGN.CENTER,
             first=True)


# --------------------------------------------------- DIAGRAM B · one signal, three
def one_signal_three(slide):
    tw, th, ty = 6.40, 1.70, 2.62
    tx = (20.0 - tw) / 2
    node(slide, tx, ty, tw, th, "UNCERTAINTY SCORE", ["How sure is the system?"],
         fill=DEEP, fg=WHITE, title_pt=16, body_pt=12)

    cy, ch = 7.05, 2.05
    cells = _row(ML, ML + CW, 3, 1.90)
    cspec = [("PERFORMANCE", ["Flag or verify the answer"]),
             ("COST", ["Should have used a cheaper model"]),
             ("RESPONSIBILITY", ["Escalate for human review"])]
    labels = [["low confidence", "= likely wrong"],
              ["low confidence", "= overspent compute"],
              ["low confidence", "= needs human eyes"]]
    lw = 2.9
    for (x, w), (title, body), lab in zip(cells, cspec, labels):
        node(slide, x, cy, w, ch, title, body, fill=ACCENT, fg=WHITE, title_pt=16,
             body_pt=12)
        x_mid = x + w / 2
        side = 1 if x_mid > 10.0 else (-1 if x_mid < 10.0 else 0)
        # start on the underside of the parent, fan out to each child
        sx = tx + tw * (0.5 + 0.34 * side)
        arrow(slide, sx, ty + th, x_mid, cy - 0.04)
        if side == 0:
            # the centre connector is vertical, so a centred label would be struck
            # through by its own line -- sit the text clear to the left of it
            edge_label(slide, x_mid - 0.16 - lw / 2, (ty + th + cy) / 2, "", lines=lab,
                       w=lw, align=PP_ALIGN.RIGHT)
        else:
            edge_label(slide, (sx + x_mid) / 2 - side * 1.25, (ty + th + cy) / 2, "",
                       lines=lab, w=lw)

    tb = textbox(slide, ML, 9.38, CW, 0.40)
    para(tb.text_frame,
         "Three planes, one signal, one converged verdict: the layer no competing "
         "product assembles.", 12, color=BLACK, first=True)


# ------------------------------------------------------- DIAGRAM C · four verdicts
def four_verdicts(slide):
    ew, eh, ey = 7.20, 1.62, 2.58
    ex = (20.0 - ew) / 2
    node(slide, ex, ey, ew, eh, "POLICY ENGINE",
         ["signals × policy × severity order"], fill=DEEP, fg=WHITE,
         title_pt=16, body_pt=12)

    vy, vh = 6.30, 2.45
    cells = _row(ML, ML + CW, 4, 0.44)
    vspec = [
        ("PASS", ["Delivered unchanged"], TINT1, DEEP, "no signal above threshold"),
        ("EDIT", ["Span redacted or claim", "softened, rest delivered"], ACCENT, WHITE,
         "transformable breach"),
        ("BLOCK", ["Policy fallback message", "instead of model text"], DEEP, WHITE,
         "non-transformable breach"),
        ("ESCALATE", ["Quarantined for human", "review, nothing delivered"], ACCENT,
         WHITE, "ambiguous or high-stakes"),
    ]
    for (x, w), (title, body, fill, fg, lab) in zip(cells, vspec):
        node(slide, x, vy, w, vh, title, body, fill=fill, fg=fg, title_pt=16,
             body_pt=11.5, stroke=RULE if fill is TINT1 else None)
        x_mid = x + w / 2
        frac = 0.5 + (x_mid - 10.0) / (CW * 1.28)
        sx = ex + ew * min(max(frac, 0.06), 0.94)
        arrow(slide, sx, ey + eh, x_mid, vy - 0.04)
        # each label sits outward of its own connector, clear of the line
        _fan_label(slide, sx, ey + eh, x_mid, vy - 0.04, (ey + eh + vy) / 2, [lab],
                   side=-1 if x_mid < 10 else 1, w=2.9)


# ------------------------------------------------ DIAGRAM 2 · three-verdict fork
def three_verdict_fork(slide):
    px, pw, ph = ML, 4.55, 2.55
    py = 5.10
    sp = rect(slide, px, py, pw, ph, GREY_BG, stroke=RULE)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.22)
    para(tf, "ONE RESPONSE", 10.5, bold=True, color=MUTED, first=True)
    para(tf, "An assistant reply that repeats a customer's SSN.", 13, bold=True,
         color=BLACK, space_before=5, line=1.06)
    para(tf, "Identical detector stack. Identical content.", 11.5, color=BLACK,
         space_before=6, line=1.05)
    para(tf, "Only the YAML differs.", 11.5, bold=True, color=ACCENT, space_before=2)

    rows = [
        ("support_bot", "Customer-facing support", "redact_and_deliver", "EDIT",
         "Span redacted, customer keeps the answer", ACCENT),
        ("hr_copilot", "Internal HR assistant", "block_on_personal_data", "BLOCK",
         "Policy fallback delivered instead", DEEP),
        ("finance_advisor", "Regulated advisory", "escalate_high_stakes", "ESCALATE",
         "HTTP 202 · quarantined for human review", ACCENT),
    ]
    cx, cw = 6.55, 7.05
    bx, bw = 14.30, 4.94
    ys = [2.62, 5.10, 7.58]
    rh = 2.10
    for (name, use, rule, verdict, effect, vfill), ry in zip(rows, ys):
        card = rect(slide, cx, ry, cw, rh, WHITE, stroke=RULE)
        ctf = card.text_frame
        ctf.margin_left = ctf.margin_right = Inches(0.20)
        para(ctf, name + ".yaml", 14, bold=True, color=DEEP, first=True)
        para(ctf, use, 11, color=MUTED, space_before=3)
        para(ctf, "policy: " + rule, 11, color=BLACK, space_before=5)
        badge = rect(slide, bx, ry + 0.20, bw, rh - 0.40, vfill,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
        btf = badge.text_frame
        para(btf, verdict, 17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True)
        para(btf, effect, 11, color=WHITE, align=PP_ALIGN.CENTER, space_before=4,
             line=1.02)
        arrow(slide, px + pw + 0.06, py + ph / 2, cx - 0.06, ry + rh / 2)
        arrow(slide, cx + cw + 0.06, ry + rh / 2, bx - 0.06, ry + rh / 2)


# ---------------------------------------------------- DIAGRAM A · escalation ladder
def escalation_ladder(slide):
    y, h = 3.86, 3.60
    cells = _row(ML, ML + CW, 5, 0.40)
    spec = [
        ("Rung 0", ["Deterministic checks", "PII · Injection · Policy"],
         "All traffic · sub-millisecond", TINT1, DEEP),
        ("Rung 1", ["Small model answers"], "All traffic · fast", TINT1, DEEP),
        ("Rung 2", ["Grounding check", "Claim vs your documents"],
         "All traffic · within hold budget", ACCENT, WHITE),
        ("Rung 3", ["Frontier model"], "Escalated traffic only", ACCENT, WHITE),
        ("Rung 4", ["Human review"], "Quarantined only", DEEP, WHITE),
    ]
    for (x, w), (title, body, tail, fill, fg) in zip(cells, spec):
        node(slide, x, y, w, h, title, body, fill=fill, fg=fg, title_pt=15, body_pt=11.5,
             italic_tail=tail, stroke=RULE if fill is TINT1 else None)
    for i in range(4):
        arrow(slide, cells[i][0] + cells[i][1] + 0.04, y + h / 2,
              cells[i + 1][0] - 0.04, y + h / 2)
    cap = textbox(slide, ML, y + h + 0.48, CW, 0.60)
    para(cap.text_frame,
         "Escalation is governed by blast radius: how irreversible the next action "
         "is. A brainstorm ships at Rung 1; an auto-refund waits at Rung 4.", 11,
         color=MUTED, align=PP_ALIGN.CENTER, first=True, line=1.06)


# ------------------------------------------------------- DIAGRAM E · market funnel
def market_funnel(slide):
    bars = [
        ("TAM", "$5.78B by 2029", "Broad AI governance + MLOps/LLMOps + privacy software",
         "45.3% CAGR · MarketsandMarkets [M-04]", 11.00, DEEP, WHITE),
        ("SAM", "$1.42B by 2030",
         "Pure-play governance software and third-party auditing, our planning basis",
         "35.7% CAGR · Grand View Research [M-04]", 2.70, ACCENT, WHITE),
        ("SOM", "$5–8M ARR",
         "Year-3 target: 120–180 accounts at blended ACV",
         "0.35–0.56% of the 2030 pure-play estimate [model]", 1.15, TINT1, DEEP),
    ]
    y = 2.80
    for i, (tag, figure, scope, src, w, fill, fg) in enumerate(bars):
        h = 1.16
        label = f"{tag}   {figure}"
        sp = rect(slide, ML, y, w, h, fill)
        if text_width(label, 17 * 1.5, True) + 0.48 <= w:
            tf = sp.text_frame
            tf.margin_left = Inches(0.24)
            para(tf, label, 17, bold=True, color=fg, align=PP_ALIGN.LEFT, first=True)
        else:  # bar too narrow for its own label -- set the label beside the bar
            ob = textbox(slide, ML + w + 0.22, y, 6.0, h, anchor=MSO_ANCHOR.MIDDLE)
            para(ob.text_frame, label, 17, bold=True,
                 color=DEEP if fill is TINT1 else fill, first=True)
        tb = textbox(slide, ML, y + h + 0.14, 11.00, 0.60)
        para(tb.text_frame, scope, 12, color=BLACK, first=True, line=1.04)
        para(tb.text_frame, src, 10.5, color=MUTED, space_before=2)
        y += h + 0.98
        if i == 2:
            nb = textbox(slide, ML, y + 0.10, 11.00, 0.34)
            para(nb.text_frame,
                 "SOM bar shown at minimum visible width; TAM and SAM bars are to scale.",
                 10, italic=True, color=MUTED, first=True)

    box = rect(slide, 12.05, 2.80, 7.19, 3.44, GREY_BG, stroke=RULE)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.20)
    para(tf, "WHY THE TWO ESTIMATES DIFFER 4.08×", 10.5, bold=True, color=ACCENT,
         first=True)
    para(tf, "Definitional, not a disagreement about growth. The broader figure absorbs "
             "adjacent developer tooling; the narrower one counts pure-play governance "
             "only.", 11.5, color=BLACK, space_before=6, line=1.06)
    para(tf, "We quote both and plan on the narrower. Both agree on a mid-thirties-percent-"
             "plus CAGR through the end of the decade.", 11.5, color=BLACK,
         space_before=6, line=1.06)


# --------------------------------------------------- DIAGRAM D · positioning map
def positioning_map(slide):
    x0, x1 = 3.55, 17.15
    y0, y1 = 2.72, 8.62
    rect(slide, x0, y0, x1 - x0, y1 - y0, WHITE, stroke=RULE)
    rect(slide, (x0 + x1) / 2, y0, (x1 - x0) / 2, (y1 - y0) / 2, GREY_BG)
    hline(slide, x0, (y0 + y1) / 2, x1 - x0, color=RULE, width=1.0)
    arrow(slide, (x0 + x1) / 2, y0, (x0 + x1) / 2, y1, color=RULE, width=1.0, head=False)

    xa = textbox(slide, x0, y1 + 0.30, x1 - x0, 0.34)
    para(xa.text_frame, "Post-hoc observation          →          Pre-delivery decision",
         11.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER, first=True)
    ya = textbox(slide, 0.80, y0, 2.55, y1 - y0, anchor=MSO_ANCHOR.MIDDLE)
    para(ya.text_frame, "Converged governance\n↑\nSingle-plane point tool", 11.5,
         bold=True, color=BLACK, align=PP_ALIGN.RIGHT, first=True, line=1.30)

    dots = [
        ("Datadog", 0.20, 0.86), ("Arize", 0.11, 0.74), ("Langfuse", 0.09, 0.60),
        ("LangSmith", 0.15, 0.42), ("LiteLLM", 0.52, 0.22), ("Portkey", 0.63, 0.42),
        ("Kong AI Gateway", 0.70, 0.30), ("NeMo Guardrails", 0.66, 0.10),
        ("Bedrock Guardrails", 0.79, 0.17), ("Azure Content Safety", 0.90, 0.27),
        ("Lakera", 0.95, 0.40), ("Presidio", 0.38, 0.11), ("Google DLP", 0.30, 0.04),
    ]
    d = 0.17
    for name, fx, fy in dots:
        cx = x0 + fx * (x1 - x0)
        cy = y1 - fy * (y1 - y0)
        rect(slide, cx - d / 2, cy - d / 2, d, d, MUTED, shape=MSO_SHAPE.OVAL)
        right = fx < 0.60
        lw = text_width(name, 10.5 * 1.5) + 0.06  # sized to its own text: no phantom collisions
        lx = cx + d / 2 + 0.10 if right else cx - d / 2 - 0.10 - lw
        tb = textbox(slide, lx, cy - 0.15, lw, 0.30, anchor=MSO_ANCHOR.MIDDLE)
        para(tb.text_frame, name, 10.5, color=BLACK,
             align=PP_ALIGN.LEFT if right else PP_ALIGN.RIGHT, first=True)

    sx = x0 + 0.845 * (x1 - x0)
    sy = y1 - 0.845 * (y1 - y0)
    rect(slide, sx - 0.26, sy - 0.26, 0.52, 0.52, ACCENT, shape=MSO_SHAPE.STAR_5_POINT)
    lb = textbox(slide, sx - 4.30, sy - 0.34, 3.95, 0.68, anchor=MSO_ANCHOR.MIDDLE)
    para(lb.text_frame, "ControlPlane", 14, bold=True, color=ACCENT,
         align=PP_ALIGN.RIGHT, first=True)
    para(lb.text_frame, "converged, pre-delivery", 10.5, color=MUTED,
         align=PP_ALIGN.RIGHT)

    qb = textbox(slide, (x0 + x1) / 2 + 0.24, y0 + 0.20, 5.4, 0.34)
    para(qb.text_frame, "THE EMPTY QUADRANT", 10.5, bold=True, color=ACCENT, first=True)


# -------------------------------------------------------- DIAGRAM F · roadmap line
def roadmap_timeline(slide):
    y = 5.85
    hline(slide, ML + 0.30, y, CW - 0.60, color=RULE, width=2.0)
    stops = [
        ("P0 · Prototype", "DONE: this submission",
         ["Gateway, six live detectors", "Policy engine, audit + HITL",
          "Console, demo, evidence"], DEEP, True),
        ("P1 · Pilot-ready", "0–3 months",
         ["Cost-plane enforcement", "Indic PII locales; scorer",
          "Presidio head-to-head", "Partner-traffic tuning"], ACCENT, False),
        ("P2 · Product", "3–9 months",
         ["Console GA; policy packs", "Compiled data plane",
          "Multi-tenant isolation; SOC 2"], ACCENT, False),
        ("P3 · Scale", "9–18 months",
         ["K8s sidecar and service mesh", "Agent-action gating",
          "Marketplace listings; SI packs"], DEEP, False),
    ]
    cells = _row(ML + 0.30, ML + CW - 0.30, 4, 0.0)
    for i, ((x, w), (title, when, items, fill, done)) in enumerate(zip(cells, stops)):
        cx = x + w / 2
        rect(slide, cx - 0.16, y - 0.16, 0.32, 0.32, fill, shape=MSO_SHAPE.OVAL)
        above = i % 2 == 0
        bh = 2.35
        by = y - 0.46 - bh if above else y + 0.46
        card = rect(slide, cx - w / 2 + 0.22, by, w - 0.44, bh, WHITE, stroke=RULE)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.14)
        para(tf, title, 14, bold=True, color=DEEP, first=True)
        para(tf, when, 10.5, bold=True, color=ACCENT if not done else DEEP,
             space_before=3)
        for it in items:
            para(tf, "–  " + it, 11, color=BLACK, space_before=4, line=1.02)
        arrow(slide, cx, by + bh if above else by, cx, y - 0.18 if above else y + 0.18,
              color=RULE, width=1.0, head=False)

    cap = textbox(slide, ML, 9.50, CW, 0.40)
    para(cap.text_frame,
         "Our roadmap is our limitations register: each P1 item closes a numbered "
         "standing limitation published in the repository.", 12, color=BLACK, first=True)


# --------------------------------------------------------- DIAGRAM H · use of funds
def use_of_funds(slide):
    y, h = 3.55, 1.55
    segs = [
        (0.60, "60%", "Founding team + one senior engineer", DEEP, WHITE),
        (0.20, "20%", "Pilot infrastructure and compliance advisory", ACCENT, WHITE),
        (0.20, "20%", "GTM, developer relations, open-source community", TINT1, DEEP),
    ]
    x = ML
    for frac, pct, label, fill, fg in segs:
        w = CW * frac
        sp = rect(slide, x, y, w, h, fill)
        para(sp.text_frame, pct, 24, bold=True, color=fg, align=PP_ALIGN.CENTER,
             first=True)
        tb = textbox(slide, x, y + h + 0.26, w, 0.90)
        para(tb.text_frame, label, 11.5, color=BLACK, align=PP_ALIGN.CENTER, first=True,
             line=1.06)
        x += w
    for frac_x in (0.60, 0.80):
        gx = ML + CW * frac_x
        arrow(slide, gx, y, gx, y + h, color=WHITE, width=1.5, head=False)
