"""Render the built deck to PNGs and report layout defects.

There is no LibreOffice on this machine, so this walks the saved .pptx shape tree and
emits SVG (rendered by headless chromium). Text is measured with Liberation Sans, which
is metric-compatible with Arial, so wrap points and overflow verdicts are truthful
rather than approximate.

Usage:
  python scripts/render_check.py                 # all slides
  python scripts/render_check.py 8 11 12         # only these slide numbers

Target deck and output directory are overridable, because there are now two decks in
the repo and the older one must stay renderable for comparison:
  CP_DECK=<abs .pptx>  CP_RENDER_OUT=<abs dir>  python scripts/render_check.py
"""

from __future__ import annotations

import base64
import html
import os
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
DECK = Path(os.environ.get("CP_DECK") or ROOT / "submission" / "b24bb1029_ControlPlane_Business_Proposal.pptx")
OUTDIR = Path(os.environ.get("CP_RENDER_OUT") or "/tmp/deck_render")
PPI = 80.0  # render px per inch
EMU = 914400.0

sys.path.insert(0, str(Path(__file__).parent))
from deck_lib import text_width as measure  # single source of truth for metrics


# ------------------------------------------------------------------ shape colour probes
def _srgb(el):
    if el is None:
        return None
    c = el.find(qn("a:srgbClr"))
    return c.get("val") if c is not None else None


def slide_bg(slide):
    """Slide-level background colour, or None.

    The colour sits at p:bg/p:bgPr/a:solidFill/a:srgbClr, several levels below p:bg, so the
    direct-child lookup in _srgb() cannot reach it. Ignoring this painted every ink-background
    slide white and made white body text look like a contrast failure.
    """
    bg = slide._element.find(f'.//{qn("p:bg")}')   # nests under p:cSld, not p:sld
    if bg is None:
        return None
    c = bg.find(f'.//{qn("a:srgbClr")}')
    return c.get("val") if c is not None else None


def has_blip_fill(shape):
    """True if this shape carries a picture fill."""
    sp = shape._element.find(qn("p:spPr"))
    return sp is not None and sp.find(qn("a:blipFill")) is not None


def blip_href(shape):
    """Inline data URI for a shape's image, or None.

    The locked template's cover, team, video and thank-you slides are picture-filled
    freeforms carrying white text. Leaving them unpainted rendered those slides as blank
    white pages with invisible type, which is indistinguishable from a real contrast defect.
    """
    el = shape._element
    blip = el.find(f'.//{qn("a:blip")}')
    if blip is None:
        return None
    rid = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
    if not rid:
        return None
    try:
        part = shape.part.related_part(rid)
        return f"data:{part.content_type};base64," + base64.b64encode(part.blob).decode()
    except Exception:
        return None


def fill_of(shape):
    spPr = getattr(shape, "_element", None)
    if spPr is None:
        return None
    sp = shape._element.find(qn("p:spPr"))
    if sp is None:
        return None
    if sp.find(qn("a:noFill")) is not None:
        return None
    return _srgb(sp.find(qn("a:solidFill")))


def line_of(shape):
    sp = shape._element.find(qn("p:spPr"))
    if sp is None:
        return None, 0.0
    ln = sp.find(qn("a:ln"))
    if ln is None:
        return None, 0.0
    if ln.find(qn("a:noFill")) is not None:
        return None, 0.0
    w = float(ln.get("w") or 12700) / 12700.0
    return _srgb(ln.find(qn("a:solidFill"))), w


def has_arrowhead(shape):
    sp = shape._element.find(qn("p:spPr"))
    if sp is None:
        return False
    ln = sp.find(qn("a:ln"))
    return ln is not None and ln.find(qn("a:tailEnd")) is not None


def prst_of(shape):
    sp = shape._element.find(qn("p:spPr"))
    if sp is None:
        return "rect"
    g = sp.find(qn("a:prstGeom"))
    return g.get("prst") if g is not None else "custom"


# ----------------------------------------------------------------------- text layout
def para_lines(tf, avail_in: float, defaults=(12.0, False, False, "000000"), fscale=1.0):
    """Lay out a text frame -> (lines, total_height_in, max_line_width_in).

    lines: list of dicts {segs, height_in, align, space_before_in, width_in}
    """
    out = []
    for p in tf.paragraphs:
        segs = []
        runs = {id(r._r): r for r in p.runs}
        for child in p._p:
            if child.tag == qn("a:br"):
                segs.append({"br": True})
                continue
            r = runs.get(id(child))
            if r is None:
                continue
            sz = (r.font.size.pt if r.font.size else defaults[0]) * fscale
            bold = bool(r.font.bold)
            italic = bool(r.font.italic)
            try:
                col = str(r.font.color.rgb) if r.font.color and r.font.color.type is not None else defaults[3]
            except Exception:
                col = defaults[3]
            segs.append({"t": r.text, "pt": sz, "b": bold, "i": italic, "c": col})
        if not any("t" in sg for sg in segs):
            segs = [{"t": "", "pt": defaults[0], "b": False, "i": False, "c": defaults[3]}]
        align = p.alignment or PP_ALIGN.LEFT
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        sb = (p.space_before.pt if p.space_before else 0.0) / 72.0
        max_pt = max((sg["pt"] for sg in segs if "t" in sg), default=defaults[0])

        # wrap into visual lines, carrying run boundaries
        cur, cur_w, lines = [], 0.0, []
        for sg in segs:
            if sg.get("br"):
                lines.append((cur, cur_w))
                cur, cur_w = [], 0.0
                continue
            words = sg["t"].split(" ") if sg["t"] else [""]
            for wi, word in enumerate(words):
                piece = (" " if wi > 0 else "") + word
                if not piece and wi == 0:
                    continue
                w_in = measure(piece, sg["pt"], sg["b"], sg["i"])
                if cur and cur_w + w_in > avail_in + 1e-9:
                    lines.append((cur, cur_w))
                    piece = word
                    w_in = measure(piece, sg["pt"], sg["b"], sg["i"])
                    cur, cur_w = [], 0.0
                cur.append({**sg, "t": piece})
                cur_w += w_in
        lines.append((cur, cur_w))
        for li, (segl, w) in enumerate(lines):
            lh = max((sg["pt"] for sg in segl if "t" in sg), default=max_pt) * 1.20 * ls / 72.0
            out.append({"segs": segl, "h": lh, "align": align,
                        "sb": sb if li == 0 else 0.0, "w": w})
    total = sum(l["h"] + l["sb"] for l in out)
    widest = max((l["w"] for l in out), default=0.0)
    return out, total, widest


def emit_text(tf, x, y, w, h, svg, defects, label, *, is_table_cell=False, fscale=1.0):
    ml = (tf.margin_left or 0) / EMU
    mr = (tf.margin_right or 0) / EMU
    mt = (tf.margin_top or 0) / EMU
    mb = (tf.margin_bottom or 0) / EMU
    avail = max(w - ml - mr, 0.05)
    lines, total, widest = para_lines(tf, avail, fscale=fscale)
    if not any(sg["t"].strip() for l in lines for sg in l["segs"]):
        return
    anchor = tf.vertical_anchor
    box_h = h - mt - mb
    if anchor == MSO_ANCHOR.MIDDLE:
        ty = y + mt + max((box_h - total) / 2.0, 0.0)
    elif anchor == MSO_ANCHOR.BOTTOM:
        ty = y + mt + max(box_h - total, 0.0)
    else:
        ty = y + mt
    snippet = " ".join(sg["t"] for l in lines for sg in l["segs"] if sg.get("t")).strip()
    snippet = (snippet[:52] + "…") if len(snippet) > 52 else snippet
    if total > box_h + 0.035:
        defects.append(f"{label}: overflow {(total - box_h) * 72:.0f}pt "
                       f"(text {total:.2f} vs box {box_h:.2f}in) :: {snippet!r}")
    if widest > avail + 0.02:
        defects.append(f"{label}: unwrappable {widest:.2f} > avail {avail:.2f}in "
                       f":: {snippet!r}")

    cy = ty
    for ln in lines:
        cy += ln["sb"]
        base = cy + ln["h"] * 0.78
        lw = ln["w"]
        if ln["align"] == PP_ALIGN.CENTER:
            cx = x + ml + (avail - lw) / 2.0
        elif ln["align"] == PP_ALIGN.RIGHT:
            cx = x + ml + (avail - lw)
        else:
            cx = x + ml
        for sg in ln["segs"]:
            if sg["t"]:
                svg.append(
                    f'<text x="{cx * PPI:.1f}" y="{base * PPI:.1f}" '
                    f'font-family="Liberation Sans, Arial" '
                    f'font-size="{sg["pt"] / 72.0 * PPI:.2f}" '
                    f'{"font-weight=\"bold\" " if sg["b"] else ""}'
                    f'{"font-style=\"italic\" " if sg["i"] else ""}'
                    f'fill="#{sg["c"]}" xml:space="preserve">{html.escape(sg["t"])}</text>')
                cx += measure(sg["t"], sg["pt"], sg["b"], sg["i"])
        cy += ln["h"]


# ------------------------------------------------------------------------- shape walk
def draw(shape, ox, oy, sx, sy, svg, defects, slide_no, depth=0, fscale=1.0):
    st = str(shape.shape_type)
    x = ox + (shape.left or 0) / EMU * sx
    y = oy + (shape.top or 0) / EMU * sy
    w = (shape.width or 0) / EMU * sx
    h = (shape.height or 0) / EMU * sy

    if shape.shape_type == 6:  # GROUP
        gx = shape._element.find(qn("p:grpSpPr")).find(qn("a:xfrm"))
        chOff = gx.find(qn("a:chOff"))
        chExt = gx.find(qn("a:chExt"))
        cox = float(chOff.get("x")) / EMU
        coy = float(chOff.get("y")) / EMU
        cw = float(chExt.get("cx")) / EMU
        ch = float(chExt.get("cy")) / EMU
        nsx = (w / cw) if cw else 1.0
        nsy = (h / ch) if ch else 1.0
        for child in shape.shapes:
            draw(child, x - cox * nsx, y - coy * nsy, nsx, nsy, svg, defects, slide_no,
                 depth + 1, fscale * (nsx + nsy) / 2.0)
        return

    name = f"s{slide_no} {shape.shape_type} {shape.name!r}"

    if shape.has_table:
        tbl = shape.table
        cy = y
        for ri, row in enumerate(tbl.rows):
            rh = row.height / EMU * sy
            cx = x
            for ci in range(len(tbl.columns)):
                cw2 = tbl.columns[ci].width / EMU * sx
                cell = tbl.cell(ri, ci)
                tcPr = cell._tc.find(qn("a:tcPr"))
                fill = _srgb(tcPr.find(qn("a:solidFill"))) if tcPr is not None else None
                if fill:
                    svg.append(f'<rect x="{cx*PPI:.1f}" y="{cy*PPI:.1f}" '
                               f'width="{cw2*PPI:.1f}" height="{rh*PPI:.1f}" fill="#{fill}"/>')
                svg.append(f'<rect x="{cx*PPI:.1f}" y="{cy*PPI:.1f}" width="{cw2*PPI:.1f}" '
                           f'height="{rh*PPI:.1f}" fill="none" stroke="#D9D9D9" stroke-width="0.6"/>')
                emit_text(cell.text_frame, cx, cy, cw2, rh, svg, defects,
                          f"{name} cell[{ri}][{ci}]", is_table_cell=True, fscale=fscale)
                cx += cw2
            cy += rh
        if cy - y > h + 0.02:
            defects.append(f"{name}: table rows total {cy-y:.2f}in exceed declared {h:.2f}in")
        if cy > 10.55:
            defects.append(f"{name}: table bottom {cy:.2f}in runs into the footer band")
        return

    if shape._element.tag.endswith("}cxnSp"):
        col, lw = line_of(shape)
        flipH = shape._element.find(qn("p:spPr")).find(qn("a:xfrm")).get("flipH") == "1"
        flipV = shape._element.find(qn("p:spPr")).find(qn("a:xfrm")).get("flipV") == "1"
        x1, x2 = (x + w, x) if flipH else (x, x + w)
        y1, y2 = (y + h, y) if flipV else (y, y + h)
        marker = ' marker-end="url(#ah)"' if has_arrowhead(shape) else ""
        svg.append(f'<line x1="{x1*PPI:.1f}" y1="{y1*PPI:.1f}" x2="{x2*PPI:.1f}" '
                   f'y2="{y2*PPI:.1f}" stroke="#{col or "808080"}" '
                   f'stroke-width="{max(lw,0.75)*PPI/72:.2f}"{marker}/>')
        return

    if shape.shape_type == 13 or has_blip_fill(shape):
        href = blip_href(shape)
        if href:
            svg.append(f'<image x="{x*PPI:.1f}" y="{y*PPI:.1f}" width="{w*PPI:.1f}" '
                       f'height="{h*PPI:.1f}" href="{href}" '
                       f'preserveAspectRatio="xMidYMid slice"/>')
            if shape.has_text_frame:
                emit_text(shape.text_frame, x, y, w, h, svg, defects, name, fscale=fscale)
            return

    fill = fill_of(shape)
    scol, slw = line_of(shape)
    prst = prst_of(shape)
    if fill or scol:
        f = f'#{fill}' if fill else "none"
        stroke = f' stroke="#{scol}" stroke-width="{max(slw,0.5)*PPI/72:.2f}"' if scol else ""
        if prst == "ellipse":
            svg.append(f'<ellipse cx="{(x+w/2)*PPI:.1f}" cy="{(y+h/2)*PPI:.1f}" '
                       f'rx="{w/2*PPI:.1f}" ry="{h/2*PPI:.1f}" fill="{f}"{stroke}/>')
        elif prst == "star5":
            cx, cy = (x + w / 2) * PPI, (y + h / 2) * PPI
            r = min(w, h) / 2 * PPI
            import math
            pts = []
            for i in range(10):
                ang = -math.pi / 2 + i * math.pi / 5
                rr = r if i % 2 == 0 else r * 0.42
                pts.append(f"{cx + rr*math.cos(ang):.1f},{cy + rr*math.sin(ang):.1f}")
            svg.append(f'<polygon points="{" ".join(pts)}" fill="{f}"{stroke}/>')
        else:
            rx = min(w, h) * 0.16 * PPI if prst == "roundRect" else 0
            svg.append(f'<rect x="{x*PPI:.1f}" y="{y*PPI:.1f}" width="{w*PPI:.1f}" '
                       f'height="{h*PPI:.1f}" rx="{rx:.1f}" fill="{f}"{stroke}/>')

    if shape.has_text_frame:
        emit_text(shape.text_frame, x, y, w, h, svg, defects, name, fscale=fscale)

    if depth == 0 and (x < -0.02 or y < -0.02 or x + w > 20.02 or y + h > 11.27):
        defects.append(f"{name}: off-slide bbox x={x:.2f} y={y:.2f} "
                       f"r={x+w:.2f} b={y+h:.2f}")


def collect_boxes(shape, ox, oy, sx, sy, out, depth=0):
    x = ox + (shape.left or 0) / EMU * sx
    y = oy + (shape.top or 0) / EMU * sy
    w = (shape.width or 0) / EMU * sx
    h = (shape.height or 0) / EMU * sy
    if shape.shape_type == 6:
        return
    if shape._element.tag.endswith("}cxnSp"):
        return
    has_text = shape.has_text_frame and shape.text_frame.text.strip()
    filled = fill_of(shape) is not None
    if not (has_text or filled):
        return
    # full-bleed backgrounds and accent spines are intentional underlays
    if w > 19.0 and h > 10.0:
        return
    if w < 0.45 and h > 8.0:
        return
    out.append({"name": f"{shape.name}", "x": x, "y": y, "w": w, "h": h,
                "text": (shape.text_frame.text.strip()[:34] if has_text else ""),
                "filled": filled})


def check_overlaps(slide, defects, slide_no):
    boxes = []
    for sh in slide.shapes:
        if sh.has_table:
            continue
        collect_boxes(sh, 0.0, 0.0, 1.0, 1.0, boxes)
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            ox = min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"])
            oy = min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"])
            if ox > 0.06 and oy > 0.06:
                # a filled card legitimately sits under its own label only if one is text-only
                area = min(a["w"] * a["h"], b["w"] * b["h"])
                # full containment is a container/underlay relationship, not a collision
                inside = ((a["x"] <= b["x"] + 0.02 and a["y"] <= b["y"] + 0.02 and
                           a["x"] + a["w"] >= b["x"] + b["w"] - 0.02 and
                           a["y"] + a["h"] >= b["y"] + b["h"] - 0.02) or
                          (b["x"] <= a["x"] + 0.02 and b["y"] <= a["y"] + 0.02 and
                           b["x"] + b["w"] >= a["x"] + a["w"] - 0.02 and
                           b["y"] + b["h"] >= a["y"] + a["h"] - 0.02))
                if inside:
                    continue
                if ox * oy < area * 0.22:
                    continue
                if a["filled"] and b["filled"] or (a["text"] and b["text"]):
                    defects.append(
                        f"s{slide_no} OVERLAP {ox:.2f}x{oy:.2f}in: "
                        f"{a['name']!r} {a['text']!r} vs {b['name']!r} {b['text']!r}")


def render(numbers=None):
    OUTDIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(DECK))
    all_defects = {}
    for i, slide in enumerate(prs.slides, 1):
        if numbers and i not in numbers:
            continue
        svg = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{20*PPI:.0f}" '
               f'height="{11.25*PPI:.0f}" viewBox="0 0 {20*PPI:.0f} {11.25*PPI:.0f}">',
               '<defs><marker id="ah" viewBox="0 0 10 10" refX="9" refY="5" '
               'markerWidth="5" markerHeight="5" orient="auto-start-reverse">'
               '<path d="M 0 1 L 10 5 L 0 9 z" fill="#808080"/></marker></defs>',
               f'<rect width="100%" height="100%" fill="#{slide_bg(slide) or "FFFFFF"}"/>']
        defects = []
        for sh in slide.shapes:
            draw(sh, 0.0, 0.0, 1.0, 1.0, svg, defects, i)
        check_overlaps(slide, defects, i)
        svg.append("</svg>")
        sp = OUTDIR / f"slide-{i:02d}.svg"
        sp.write_text("\n".join(svg))
        subprocess.run(["chromium", "--headless", "--no-sandbox", "--disable-gpu",
                        "--hide-scrollbars", "--force-device-scale-factor=1",
                        f"--window-size={int(20*PPI)},{int(11.25*PPI)}",
                        f"--screenshot={OUTDIR}/slide-{i:02d}.png", str(sp)],
                       capture_output=True, timeout=120)
        if defects:
            all_defects[i] = defects
    print(f"rendered -> {OUTDIR}")
    if all_defects:
        total = sum(len(v) for v in all_defects.values())
        print(f"\n=== {total} LAYOUT DEFECTS on {len(all_defects)} slides ===")
        for k in sorted(all_defects):
            print(f"\n-- slide {k}")
            for d in all_defects[k]:
                print("   ", d)
    else:
        print("no layout defects detected")
    return all_defects


if __name__ == "__main__":
    nums = [int(a) for a in sys.argv[1:]] or None
    render(nums)
