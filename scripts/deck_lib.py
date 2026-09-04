"""Shared drawing primitives for the ControlPlane business-proposal deck.

The brand template is 20 x 11.25 in (18288000 x 10287000 EMU) -- exactly 1.5x the
13.333 x 7.5 in grid the deck brief is written against. Every length and font size in
the brief is therefore multiplied by SCALE so proportions land where the brief intends.
"""

from __future__ import annotations

import copy
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- geometry
SCALE = 1.5  # brief grid (13.333 x 7.5) -> template grid (20 x 11.25)

SW, SH = 20.0, 11.25  # slide, inches
ML, MR = 0.76, 19.24  # content margins
CW = MR - ML  # 18.48 usable width

TITLE_Y = 0.82
RULE_Y = 1.86
BODY_TOP = 2.20
BODY_BOT = 9.70
SRC_Y = 9.86
FOOTER_Y = 10.55

# diagram safe region (brief x 0.6-12.7, y 1.6-6.6 scaled by 1.5)
DIA_X0, DIA_X1 = 0.90, 19.05
DIA_Y0, DIA_Y1 = 2.45, 9.60

FOOTER_TEXT = "Copyright © 2026 Accenture. All rights reserved."

# ---------------------------------------------------------------------------- palette
ACCENT = "A100FF"  # Accenture purple
DEEP = "460073"  # deep purple
TINT1 = "E6BFFF"
TINT2 = "ECCCFF"
TINT3 = "C1A3FF"
BLACK = "000000"
WHITE = "FFFFFF"
GREY_BG = "F2F2F2"
GREY_BG2 = "EAEAEA"
RULE = "D9D9D9"
MUTED = "808080"

FONT = "Arial"


_FONT_FILES = {
    (False, False): "/usr/share/fonts/liberation/LiberationSans-Regular.ttf",
    (True, False): "/usr/share/fonts/liberation/LiberationSans-Bold.ttf",
    (False, True): "/usr/share/fonts/liberation/LiberationSans-Italic.ttf",
    (True, True): "/usr/share/fonts/liberation/LiberationSans-BoldItalic.ttf",
}
_MEASURE_REF = 200  # measure large, then scale: keeps fractional point sizes exact
_font_cache: dict = {}


def _measure_font(bold: bool, italic: bool):
    """Liberation Sans is metric-compatible with Arial, so widths transfer."""
    key = (bold, italic)
    if key not in _font_cache:
        from PIL import ImageFont

        _font_cache[key] = ImageFont.truetype(_FONT_FILES[key], _MEASURE_REF)
    return _font_cache[key]


def text_width(text: str, size_pt: float, bold: bool = False, italic: bool = False) -> float:
    """Advance width of `text` in INCHES at an absolute point size."""
    if not text:
        return 0.0
    return _measure_font(bold, italic).getlength(text) / _MEASURE_REF * size_pt / 72.0


def wrapped_lines(text: str, avail_in: float, size_pt: float, bold: bool = False,
                  italic: bool = False) -> int:
    """How many visual lines `text` needs inside `avail_in`."""
    n, cur = 1, 0.0
    for i, word in enumerate(str(text).split(" ")):
        piece = (" " if i else "") + word
        w = text_width(piece, size_pt, bold, italic)
        if cur and cur + w > avail_in + 1e-9:
            n += 1
            cur = text_width(word, size_pt, bold, italic)
        else:
            cur += w
    return n


def fs(points: float) -> Pt:
    """Brief-relative font size -> template-absolute."""
    return Pt(round(points * SCALE, 1))


def rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


# ------------------------------------------------------------------------- primitives
def _noline(shape):
    shape.line.fill.background()
    return shape


def _solid(shape, hex6: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hex6)
    return shape


def _stroke(shape, hex6: str, width_pt: float = 0.75):
    shape.line.color.rgb = rgb(hex6)
    shape.line.width = Pt(width_pt)
    return shape


def style_run(run, size, *, bold=False, italic=False, color=BLACK, font=FONT):
    run.font.name = font
    run.font.size = size if isinstance(size, Pt) else fs(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    # east-asian / complex-script faces too, so nothing falls back to a theme font
    rpr = run.font._rPr
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", font)
    return run


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb


def para(tf, text, size, *, bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, line=None, first=False, font=FONT):
    """Append a paragraph (or fill the first, empty one)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    if text:
        # a literal "\n" inside <a:t> is NOT a line break in PowerPoint -- only <a:br/> is
        parts = str(text).split("\n")
        for i, part in enumerate(parts):
            if i:
                p._p.append(p._p.makeelement(qn("a:br"), {}))
            style_run(p.add_run(), size, bold=bold, italic=italic, color=color,
                      font=font).text = part
    else:
        # keep an empty paragraph from inheriting a theme font
        style_run(p.add_run(), size, color=color, font=font).text = ""
    return p


def rect(slide, x, y, w, h, fill=None, *, shape=MSO_SHAPE.RECTANGLE, radius=None,
         stroke=None, stroke_w=0.75):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None and shape in (MSO_SHAPE.ROUNDED_RECTANGLE,):
        sp.adjustments[0] = radius
    if fill is None:
        sp.fill.background()
    else:
        _solid(sp, fill)
    if stroke is None:
        _noline(sp)
    else:
        _stroke(sp, stroke, stroke_w)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sp


def node(slide, x, y, w, h, title, body=None, *, fill=ACCENT, fg=WHITE, radius=0.10,
         title_pt=15, body_pt=11.5, stroke=None, italic_tail=None, align=PP_ALIGN.CENTER):
    """A rounded diagram node: bold title, optional body lines, optional italic tail."""
    sp = rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius,
              stroke=stroke)
    tf = sp.text_frame
    para(tf, title, title_pt, bold=True, color=fg, align=align, first=True, line=0.92)
    for i, ln in enumerate(body or []):
        para(tf, ln, body_pt, color=fg, align=align, space_before=3 if i == 0 else 1,
             line=0.92)
    if italic_tail:
        para(tf, italic_tail, body_pt - 0.5, italic=True, color=fg, align=align,
             space_before=4, line=0.92)
    return sp


def _add_arrowhead(shape, kind="triangle"):
    ln = shape.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": kind, "w": "med", "len": "med"})
    ln.append(tail)


def arrow(slide, x1, y1, x2, y2, *, color=MUTED, width=1.5, head=True):
    cx = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    cx.line.color.rgb = rgb(color)
    cx.line.width = Pt(width)
    if head:
        _add_arrowhead(cx)
    return cx


def hline(slide, x, y, w, *, color=RULE, width=1.0):
    return arrow(slide, x, y, x + w, y, color=color, width=width, head=False)


def edge_label(slide, cx, cy, text, *, w=2.6, pt=10, color=MUTED, lines=None,
               align=PP_ALIGN.CENTER):
    """Small label centred on (cx, cy) -- used at connector midpoints.

    `align` matters when a label has to sit *beside* a vertical connector rather
    than on it: LEFT lets the caller anchor the box just clear of the line.
    """
    body = lines or [text]
    h = 0.20 * len(body) + 0.10
    tb = textbox(slide, cx - w / 2, cy - h / 2, w, h, anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(body):
        para(tb.text_frame, ln, pt, color=color, align=align, first=(i == 0),
             line=0.92)
    return tb


# ------------------------------------------------------------------------ slide chrome
@dataclass
class Deck:
    prs: object
    blank_layout: object

    def new(self):
        s = self.prs.slides.add_slide(self.blank_layout)
        # drop inherited date/footer/slide-number placeholders; we draw our own
        for ph in list(s.placeholders):
            ph._element.getparent().remove(ph._element)
        return s


def footer(slide, number, *, dark=False):
    col = RULE if dark else MUTED
    tb = textbox(slide, 14.82, FOOTER_Y, 3.93, 0.24)
    para(tb.text_frame, FOOTER_TEXT, Pt(12), color=col, first=True)
    nb = textbox(slide, 18.98, FOOTER_Y, 0.60, 0.24)
    para(nb.text_frame, str(number), Pt(12), color=col, align=PP_ALIGN.RIGHT, first=True)


def slide_title(slide, title, *, kicker=None, rule=True, pt=25):
    if kicker:
        kb = textbox(slide, ML, TITLE_Y - 0.34, CW, 0.30)
        para(kb.text_frame, kicker.upper(), 9.5, bold=True, color=ACCENT, first=True)
    tb = textbox(slide, ML, TITLE_Y, CW, 0.95)
    para(tb.text_frame, title, pt, bold=True, color=BLACK, first=True, line=0.95)
    if rule:
        r = rect(slide, ML, RULE_Y, 1.55, 0.055, ACCENT)
        r.shadow.inherit = False
    return tb


def source_line(slide, text, *, y=SRC_Y):
    tb = textbox(slide, ML, y, CW, 0.30)
    para(tb.text_frame, text, 9.5, color=MUTED, first=True)
    return tb


def lead(slide, text, y=BODY_TOP, *, w=None, pt=13.5, color=BLACK, bold=False):
    width = w or CW
    n = wrapped_lines(text, width, pt * SCALE, bold)
    h = n * (pt * SCALE) * 1.20 * 1.05 / 72.0 + 0.04
    tb = textbox(slide, ML, y, width, h)
    para(tb.text_frame, text, pt, bold=bold, color=color, first=True, line=1.05)
    return tb


def bullets(slide, x, y, w, items, *, pt=12, gap=0.40, color=BLACK, marker=ACCENT,
            marker_w=0.13):
    """Left-aligned bullets with a small square marker. items: list[str] or (bold, rest)."""
    cy = y
    for it in items:
        rect(slide, x, cy + 0.055, marker_w, marker_w, marker)
        tb = textbox(slide, x + marker_w + 0.18, cy - 0.04, w - marker_w - 0.18, gap)
        p = para(tb.text_frame, "", pt, color=color, first=True, line=1.02)
        p.runs[0].text = ""
        if isinstance(it, tuple):
            style_run(p.add_run(), pt, bold=True, color=color).text = it[0]
            style_run(p.add_run(), pt, color=color).text = it[1]
        else:
            style_run(p.add_run(), pt, color=color).text = it
        cy += gap
    return cy


def kpi_tile(slide, x, y, w, h, big, caption, *, tag=None, fill=TINT2, big_pt=27,
             cap_pt=10.5):
    sp = rect(slide, x, y, w, h, fill)
    tf = sp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.14)
    para(tf, big, big_pt, bold=True, color=DEEP, align=PP_ALIGN.LEFT, first=True, line=0.92)
    para(tf, caption, cap_pt, color=BLACK, align=PP_ALIGN.LEFT, space_before=3, line=1.0)
    if tag:
        para(tf, tag, 9, bold=True, color=MUTED, align=PP_ALIGN.LEFT, space_before=3)
    return sp


def _cell_fill(cell, hex6):
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    fill = tcPr.makeelement(qn("a:solidFill"), {})
    clr = tcPr.makeelement(qn("a:srgbClr"), {"val": hex6})
    fill.append(clr)
    # solidFill must precede the border elements in CT_TableCellProperties
    tcPr.insert(0, fill)


def _cell_borders(cell, hex6=RULE, width_pt=0.75, sides=("T", "B")):
    tcPr = cell._tc.get_or_add_tcPr()
    for side in sides:
        tag = qn(f"a:ln{side}")
        for old in tcPr.findall(tag):
            tcPr.remove(old)
        ln = tcPr.makeelement(tag, {"w": str(int(width_pt * 12700)), "cap": "flat"})
        sf = tcPr.makeelement(qn("a:solidFill"), {})
        sf.append(tcPr.makeelement(qn("a:srgbClr"), {"val": hex6}))
        ln.append(sf)
        tcPr.insert(0, ln)


def table(slide, x, y, w, rows, col_ratios, *, header=True, pt=11, header_pt=11,
          row_h=0.42, header_h=0.46, first_col_bold=False, align=None, zebra=True):
    """Brand-styled table. rows[0] is the header row when header=True."""
    nrow, ncol = len(rows), len(col_ratios)
    shp = slide.shapes.add_table(nrow, ncol, Inches(x), Inches(y), Inches(w),
                                 Inches(header_h + row_h * (nrow - 1)))
    tbl = shp.table
    tbl._tbl.tblPr.set("firstRow", "0")
    tbl._tbl.tblPr.set("bandRow", "0")
    total = sum(col_ratios)
    for i, r in enumerate(col_ratios):
        tbl.columns[i].width = Emu(int(Inches(w) * r / total))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(header_h if (header and ri == 0) else row_h)
        for ci in range(ncol):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.12)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_head = header and ri == 0
            if is_head:
                _cell_fill(cell, DEEP)
            elif zebra and (ri % 2 == 0):
                _cell_fill(cell, GREY_BG)
            else:
                _cell_fill(cell, WHITE)
            _cell_borders(cell)
            tf = cell.text_frame
            tf.word_wrap = True
            txt = row[ci] if ci < len(row) else ""
            al = PP_ALIGN.LEFT if align is None else align[ci]
            para(tf, txt, header_pt if is_head else pt,
                 bold=is_head or (first_col_bold and ci == 0 and not is_head),
                 color=WHITE if is_head else BLACK, align=al, first=True, line=1.0)
    return shp


def divider(slide, number, title, subtitle=None):
    """Full-bleed deep-purple section opener."""
    rect(slide, 0, 0, SW, SH, DEEP)
    rect(slide, 0, 0, 0.34, SH, ACCENT)
    nb = textbox(slide, 2.10, 3.55, 6.0, 2.30)
    para(nb.text_frame, number, 96, bold=True, color=ACCENT, first=True, line=0.86)
    tb = textbox(slide, 2.10, 5.75, 14.0, 1.30)
    para(tb.text_frame, title, 40, bold=True, color=WHITE, first=True, line=0.98)
    if subtitle:
        sb = textbox(slide, 2.10, 7.05, 13.0, 0.90)
        para(sb.text_frame, subtitle, 14, color=TINT1, first=True, line=1.10)


# ------------------------------------------------------------------- font normalisation
class _BlobDoc:
    """Adapter letting force_arial() treat a blob-only part like an XML document."""

    def __init__(self, part):
        self.part = part
        from lxml import etree

        self.root = etree.fromstring(part.blob)

    def iter(self, tag=None):
        return self.root.iter(tag) if tag is not None else self.root.iter()

    def flush(self):
        from lxml import etree

        self.part._blob = etree.tostring(
            self.root, xml_declaration=True, encoding="UTF-8", standalone=True)


def force_arial(prs):
    """Rewrite every typeface in the package to Arial and drop theme-font references."""
    docs = [prs.part._element]
    for m in prs.slide_masters:
        docs.append(m._element)
        for lay in m.slide_layouts:
            docs.append(lay._element)
    # Sweep every theme and notes master in the PACKAGE, not just the ones reachable
    # from a slide master: the notes master carries its own theme (theme2.xml), and
    # walking master rels alone left a stray Calibri font scheme in the shipped file.
    raw_parts = [part for part in prs.part.package.iter_parts()
                 if "/theme/" in str(part.partname)
                 or "notesMaster" in str(part.partname)]
    for s in prs.slides:
        docs.append(s._element)
        if s.has_notes_slide:
            docs.append(s.notes_slide._element)
    for part in raw_parts:
        el = getattr(part, "_element", None)
        if el is not None:
            docs.append(el)
        else:
            docs.append(_BlobDoc(part))

    # a:sym too -- it names the face used for symbol characters, and the template
    # leaves it pointing at Arimo, which is the one way a glyph could still render
    # in a non-Arial face after everything else has been normalised.
    latin, ea, cs = qn("a:latin"), qn("a:ea"), qn("a:cs")
    sym = qn("a:sym")
    for root in docs:
        for tag in (latin, ea, cs, sym):
            for el in root.iter(tag):
                el.set("typeface", FONT)
                for junk in ("panose", "pitchFamily", "charset"):
                    el.attrib.pop(junk, None)
        # +mj-lt / +mn-lt style references resolve through the theme; make them literal
        for el in list(root.iter()):
            for attr in ("typeface",):
                v = el.get(attr)
                if v and v.startswith("+"):
                    el.set(attr, FONT)
    for doc in docs:
        if isinstance(doc, _BlobDoc):
            doc.flush()


def bold_italic_from_facename(prs):
    """The template encodes weight in the face name (Arimo Bold / Arimo Bold Italics).

    Must run BEFORE force_arial(), while those names are still readable.
    """
    for s in prs.slides:
        for el in s._element.iter(qn("a:latin")):
            name = (el.get("typeface") or "").lower()
            rpr = el.getparent()
            if "bold" in name and rpr.get("b") is None:
                rpr.set("b", "1")
            if "italic" in name and rpr.get("i") is None:
                rpr.set("i", "1")


# ------------------------------------------------------------------- slide list surgery
def reorder_and_prune(prs, keep_order):
    """keep_order: list of 0-based original slide indices, in the order to keep them.

    Slides not listed are dropped from the presentation.
    """
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    keep = [ids[i] for i in keep_order]
    drop = [e for e in ids if e not in keep]
    for e in ids:
        sldIdLst.remove(e)
    for e in keep:
        sldIdLst.append(e)
    for e in drop:
        rid = e.get(qn("r:id"))
        try:
            prs.part.drop_rel(rid)
        except KeyError:
            pass


def clone_after(prs, index):
    """Not used for content slides; kept for reference."""
    raise NotImplementedError


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_text(shape, lines, *, sizes=None, colors=None, bolds=None, italics=None,
             aligns=None):
    """Replace a template text frame's content, keeping its position."""
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, ln in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        if aligns:
            p.alignment = aligns[i]
        style_run(p.add_run(),
                  sizes[i] if sizes else Pt(18),
                  bold=bolds[i] if bolds else False,
                  italic=italics[i] if italics else False,
                  color=colors[i] if colors else BLACK).text = ln
    return shape
